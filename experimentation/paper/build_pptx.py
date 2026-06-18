#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""Présentation AI Sovereign FinOps Operator — orientée 'une fonctionnalité = un problème'.
Contenu strictement aligné sur README.md, docs/features/*, docs/LIMITATIONS.md."""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# palette
DARK   = RGBColor(0x0F, 0x1B, 0x2D)
GREEN  = RGBColor(0x1E, 0xC8, 0x7A)
BLUE   = RGBColor(0x3B, 0x82, 0xF6)
RED    = RGBColor(0xE5, 0x3E, 0x3E)
ORANGE = RGBColor(0xF5, 0x9E, 0x0B)
GREY   = RGBColor(0x5B, 0x66, 0x70)
LIGHT  = RGBColor(0xF4, 0xF7, 0xFA)
REDBG  = RGBColor(0xFD, 0xEC, 0xEC)
GRNBG  = RGBColor(0xE9, 0xF8, 0xF1)
BLUBG  = RGBColor(0xEA, 0xF1, 0xFB)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
DARKTX = RGBColor(0x1A, 0x24, 0x30)
MUTE   = RGBColor(0xB8, 0xC6, 0xD4)

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height
BLANK = prs.slide_layouts[6]

def slide(): return prs.slides.add_slide(BLANK)

def rect(s, x, y, w, h, color, line=None):
    shp = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    shp.fill.solid(); shp.fill.fore_color.rgb = color
    if line is None: shp.line.fill.background()
    else: shp.line.color.rgb = line; shp.line.width = Pt(1)
    shp.shadow.inherit = False
    return shp

def rrect(s, x, y, w, h, color, line=None, lw=1.25):
    shp = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    shp.fill.solid(); shp.fill.fore_color.rgb = color
    if line is None: shp.line.fill.background()
    else: shp.line.color.rgb = line; shp.line.width = Pt(lw)
    shp.shadow.inherit = False
    return shp

def txt(s, x, y, w, h, runs, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
        space_after=4, ls=1.0):
    tb = s.shapes.add_textbox(x, y, w, h); tf = tb.text_frame
    tf.word_wrap = True; tf.vertical_anchor = anchor
    tf.margin_left = Pt(3); tf.margin_right = Pt(3)
    tf.margin_top = Pt(1); tf.margin_bottom = Pt(1)
    for i, para in enumerate(runs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align; p.space_after = Pt(space_after)
        p.space_before = Pt(0); p.line_spacing = ls
        for (t, sz, b, c, *rest) in para:
            it = rest[0] if rest else False
            r = p.add_run(); r.text = t
            r.font.size = Pt(sz); r.font.bold = b
            r.font.color.rgb = c; r.font.italic = it
            r.font.name = "Calibri"
    return tb

def title_bar(s, kicker, title, num):
    rect(s, 0, 0, SW, Inches(1.12), DARK)
    rect(s, 0, Inches(1.12), SW, Pt(4), GREEN)
    txt(s, Inches(0.55), Inches(0.14), Inches(10.5), Inches(0.3),
        [[(kicker, 12, True, GREEN)]])
    txt(s, Inches(0.55), Inches(0.42), Inches(11.2), Inches(0.62),
        [[(title, 25, True, WHITE)]])
    c = s.shapes.add_shape(MSO_SHAPE.OVAL, SW - Inches(1.0), Inches(0.31), Inches(0.5), Inches(0.5))
    c.fill.solid(); c.fill.fore_color.rgb = GREEN; c.line.fill.background(); c.shadow.inherit = False
    c.text_frame.text = str(num)
    pr = c.text_frame.paragraphs[0]; pr.alignment = PP_ALIGN.CENTER
    pr.runs[0].font.size = Pt(16); pr.runs[0].font.bold = True; pr.runs[0].font.color.rgb = DARK

def footer(s, txt_s="AI Sovereign FinOps Operator"):
    txt(s, Inches(0.55), SH - Inches(0.4), Inches(9), Inches(0.3),
        [[(txt_s, 9, False, GREY)]])

def feature_slide(num, kicker, title, problem, mechanism, impact,
                  formula=None, metric=None):
    """Gabarit standard : PROBLÈME (rouge) → CE QUE FAIT LE MOTEUR → IMPACT."""
    s = slide()
    title_bar(s, kicker, title, num)
    # bloc problème
    top = Inches(1.4)
    rrect(s, Inches(0.55), top, Inches(12.2), Inches(1.15), REDBG, line=RED, lw=1.0)
    txt(s, Inches(0.78), top + Inches(0.1), Inches(3), Inches(0.32),
        [[("LE PROBLÈME", 12, True, RED)]])
    txt(s, Inches(0.78), top + Inches(0.42), Inches(11.7), Inches(0.7),
        [[(problem, 14.5, False, DARKTX)]], ls=1.02)
    # deux colonnes : mécanisme / impact
    cy = Inches(2.75); ch = Inches(3.05)
    rrect(s, Inches(0.55), cy, Inches(6.0), ch, GRNBG, line=GREEN, lw=1.0)
    txt(s, Inches(0.78), cy + Inches(0.1), Inches(5.6), Inches(0.35),
        [[("CE QUE FAIT LE MOTEUR", 12, True, RGBColor(0x12,0x7A,0x4D))]])
    txt(s, Inches(0.78), cy + Inches(0.5), Inches(5.55), Inches(2.45),
        [[("•  " + l, 13, False, DARKTX)] for l in mechanism], space_after=7, ls=1.04)

    rrect(s, Inches(6.75), cy, Inches(6.0), ch, BLUBG, line=BLUE, lw=1.0)
    txt(s, Inches(6.98), cy + Inches(0.1), Inches(5.6), Inches(0.35),
        [[("IMPACT", 12, True, BLUE)]])
    txt(s, Inches(6.98), cy + Inches(0.5), Inches(5.55), Inches(2.45),
        [[("✓  " + l, 13, False, DARKTX)] for l in impact], space_after=7, ls=1.04)
    # bandeau bas : formule / métrique
    by = Inches(6.05)
    if formula:
        rrect(s, Inches(0.55), by, Inches(12.2), Inches(0.95), RGBColor(0x10,0x1B,0x2B))
        txt(s, Inches(0.8), by + Inches(0.08), Inches(11.7), Inches(0.8),
            [[(formula, 14, True, WHITE)]] +
            ([[(metric, 11.5, False, MUTE)]] if metric else []),
            space_after=3)
    elif metric:
        rrect(s, Inches(0.55), by, Inches(12.2), Inches(0.7), LIGHT)
        txt(s, Inches(0.8), by + Inches(0.1), Inches(11.7), Inches(0.5),
            [[("Métrique / sortie :  ", 12, True, DARK), (metric, 12, False, DARKTX)]],
            anchor=MSO_ANCHOR.MIDDLE)
    footer(s)
    return s

# ================================================================ 1 — Titre
s = slide()
rect(s, 0, 0, SW, SH, DARK)
rect(s, 0, 0, Inches(0.25), SH, GREEN)
rect(s, 0, Inches(4.5), SW, Pt(5), GREEN)
txt(s, Inches(0.9), Inches(1.4), Inches(11.5), Inches(0.4),
    [[("OPÉRATEUR KUBERNETES  ·  FINOPS & SOUVERAINETÉ POUR L'IA D'ENTREPRISE", 14, True, GREEN)]])
txt(s, Inches(0.9), Inches(2.0), Inches(11.6), Inches(1.6),
    [[("AI Sovereign FinOps", 50, True, WHITE)],
     [("Operator", 50, True, WHITE)]])
txt(s, Inches(0.9), Inches(4.7), Inches(11.6), Inches(1.5),
    [[("Gouverner les appels LLM d'une organisation — coût, attribution, fournisseurs,", 16, False, MUTE)],
     [("résidence des données, budgets, arbitrage API managée vs auto-hébergement — en pur déclaratif (CRDs).", 16, False, MUTE)]])
txt(s, Inches(0.9), Inches(6.4), Inches(11.6), Inches(0.5),
    [[("Catalogue out-of-the-box · 2 plans de souveraineté (gateway + eBPF/shadow-AI) · enforce réel · 7 CRDs · 16 métriques · CNCF · Apache 2.0", 13, True, GREEN)]])

# ================================================================ 2 — Les 3 problèmes
s = slide()
title_bar(s, "POURQUOI", "Trois angles morts quand l'IA arrive plus vite que sa gouvernance", 2)
txt(s, Inches(0.55), Inches(1.35), Inches(12.2), Inches(0.5),
    [[("Récurrents en secteur régulé (santé, finance, assurance, public, industrie critique) :", 14.5, False, DARKTX)]])
probs = [
    ("1 · Opacité des coûts", RED,
     ["Personne ne sait quelle", "équipe / app dépense quoi,", "ni quand un modèle moins", "cher suffirait."]),
    ("2 · Souveraineté & conformité", ORANGE,
     ["La résidence des données et", "l'usage de données sensibles", "sont gérés au cas par cas,", "voire pas du tout."]),
    ("3 · Managé vs auto-hébergé", BLUE,
     ["Aucune visibilité sur le", "volume à partir duquel", "l'auto-hébergement GPU", "devient rentable."]),
]
cx = Inches(0.55); cw = Inches(3.95); gap = Inches(0.2)
for i,(t,col,lines) in enumerate(probs):
    x = cx + i*(cw+gap)
    rrect(s, x, Inches(2.0), cw, Inches(3.0), WHITE, line=col, lw=1.5)
    rect(s, x, Inches(2.0), cw, Inches(0.6), col)
    txt(s, x+Inches(0.2), Inches(2.08), cw-Inches(0.4), Inches(0.45),
        [[(t, 15, True, WHITE)]], anchor=MSO_ANCHOR.MIDDLE)
    txt(s, x+Inches(0.25), Inches(2.8), cw-Inches(0.5), Inches(2.0),
        [[(l, 14, False, DARKTX)] for l in lines], space_after=6, ls=1.08)
rrect(s, Inches(0.55), Inches(5.25), Inches(12.2), Inches(1.45), LIGHT)
txt(s, Inches(0.78), Inches(5.4), Inches(11.8), Inches(1.2),
    [[("La réponse : ", 15, True, DARK),
      ("transformer ces décisions ad hoc en politiques Kubernetes versionnées, auditables et réconciliées en continu.", 15, False, DARKTX)],
     [("⚠️  Le produit ne promet PAS la conformité juridique : il produit une traçabilité exploitable et un", 13.5, True, RED)],
     [("dossier d'audit. Modes gradués : reportOnly (observe) → warn (alerte) → enforce (réécrit la route Envoy).", 13.5, False, DARKTX)]])
footer(s)

# ================================================================ 3 — Architecture
s = slide()
title_bar(s, "COMMENT", "L'opérateur = plan de contrôle ; la gateway reste le plan de données", 3)
rrect(s, Inches(0.55), Inches(1.45), Inches(12.2), Inches(1.5), BLUBG)
txt(s, Inches(0.75), Inches(1.52), Inches(8.5), Inches(0.32),
    [[("PLAN DE DONNÉES — observé ; en mode enforce, la route Envoy (AIGatewayRoute) est réécrite", 12, True, BLUE)]])
boxes = [("Applications", BLUE), ("Gateway IA\nEnvoy AI Gateway", BLUE), ("Fournisseurs LLM\nOpenAI · Mistral · auto-hébergé", GREY)]
bx = Inches(0.95)
for i,(t,c) in enumerate(boxes):
    w = Inches(3.4); x = bx + i*Inches(3.9)
    rrect(s, x, Inches(1.95), w, Inches(0.85), WHITE, line=c)
    txt(s, x, Inches(2.0), w, Inches(0.75),
        [[(ln, 12, True, DARKTX)] for ln in t.split("\n")],
        align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, space_after=0)
    if i < 2:
        txt(s, x+w-Inches(0.1), Inches(2.0), Inches(0.6), Inches(0.6),
            [[("→", 22, True, GREEN)]], align=PP_ALIGN.CENTER)
txt(s, Inches(3.5), Inches(3.0), Inches(8), Inches(0.35),
    [[("↓  télémétrie (Prometheus / ConfigMap / Envoy)    ↑  enforce : réécriture de la route", 12, True, GREEN)]], align=PP_ALIGN.CENTER)
rrect(s, Inches(0.55), Inches(3.45), Inches(12.2), Inches(2.05), DARK)
txt(s, Inches(0.75), Inches(3.52), Inches(9), Inches(0.32),
    [[("PLAN DE CONTRÔLE — l'opérateur, piloté par des CRDs déclaratives", 12, True, GREEN)]])
engines = ["cost\nengine", "budget\nengine", "sovereignty\nengine", "breakeven\nengine", "reco\nengine", "reporting"]
ex = Inches(0.85)
for i,e in enumerate(engines):
    w = Inches(1.82); x = ex + i*Inches(1.98)
    rrect(s, x, Inches(3.95), w, Inches(1.2), RGBColor(0x1B,0x2C,0x42), line=GREEN)
    txt(s, x, Inches(4.05), w, Inches(1.0),
        [[(ln, 12, True, WHITE)] for ln in e.split("\n")],
        align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, space_after=0)
outs = [("Métriques ai_finops_*\nPrometheus → Grafana", GREEN),
        (".status des CRDs\ncoût · budget · findings · reco", BLUE),
        ("Rapport ConfigMap\nMarkdown / JSON", ORANGE)]
ox = Inches(0.95)
for i,(t,c) in enumerate(outs):
    w = Inches(3.7); x = ox + i*Inches(3.95)
    rrect(s, x, Inches(5.7), w, Inches(0.95), WHITE, line=c)
    txt(s, x, Inches(5.75), w, Inches(0.85),
        [[(ln, 12, True, DARKTX)] for ln in t.split("\n")],
        align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, space_after=0)
footer(s)

# ================================================================ 4 — Mapping fonctionnalité ↔ problème
s = slide()
title_bar(s, "FIL CONDUCTEUR", "Chaque fonctionnalité répond à un problème précis", 4)
rows = [
    ("Collecteurs de télémétrie", "Les moteurs ont besoin de données d'usage sans dépendre d'une gateway donnée", GREY),
    ("Moteur de coût", "Opacité : qui dépense quoi, quel modèle coûte le plus", RED),
    ("Moteur de budget", "Aucune alerte graduée avant d'atteindre le budget", ORANGE),
    ("Moteur de souveraineté", "Où part la donnée ? quel fournisseur est autorisé ?", BLUE),
    ("Moteur break-even", "API managée ou GPU auto-hébergé : à partir de quand ?", GREEN),
    ("Moteur de recommandation", "Le constat ne dit pas QUOI faire, ni combien on gagne", RGBColor(0x2E,0xA0,0x6A)),
    ("Reporting", "Besoin d'un dossier d'audit lisible et reproductible", RGBColor(0x7A,0x4D,0xC8)),
]
y0 = Inches(1.45); rh = Inches(0.74)
for i,(feat,prob,col) in enumerate(rows):
    y = y0 + i*rh
    bg = LIGHT if i%2==0 else WHITE
    rect(s, Inches(0.55), y, Inches(12.2), rh-Inches(0.07), bg)
    rect(s, Inches(0.55), y, Pt(5), rh-Inches(0.07), col)
    txt(s, Inches(0.8), y, Inches(4.0), rh-Inches(0.07),
        [[(feat, 15, True, DARK)]], anchor=MSO_ANCHOR.MIDDLE)
    txt(s, Inches(4.9), y, Inches(0.6), rh-Inches(0.07),
        [[("→", 18, True, col)]], anchor=MSO_ANCHOR.MIDDLE)
    txt(s, Inches(5.5), y, Inches(7.1), rh-Inches(0.07),
        [[(prob, 13.5, False, DARKTX)]], anchor=MSO_ANCHOR.MIDDLE)
footer(s)

# ================================================================ 5 — Collecteurs
feature_slide(
    5, "FONCTIONNALITÉ — PRÉREQUIS", "Collecteurs de télémétrie : brancher en lecture seule",
    "Les moteurs ont besoin du nombre de requêtes, des tokens in/out, du modèle, du fournisseur, "
    "du namespace/app/équipe — mais ne doivent dépendre d'aucune gateway en particulier.",
    ["Abstraction TelemetryCollector : une interface, 4 implémentations.",
     "prometheus — scrute un endpoint, parse les compteurs ai_finops_*.",
     "configmap — lit des échantillons d'usage depuis un ConfigMap.",
     "aigw — Envoy AI Gateway (histogramme gen_ai_*), source réelle.",
     "fake — opt-in EXPLICITE (mode:fake) ; jamais de fallback silencieux."],
    ["Branchement non-intrusif sur une gateway existante.",
     "Pas de source réelle → erreur + condition NoTelemetrySource.",
     "Aucun chiffre inventé : crédibilité du plan de contrôle.",
     "Normalise tout en UsageSample, base commune des 6 moteurs."],
    metric="UsageSample{Namespace, Application, Team, Provider, Model, Requests, In/OutputTokens, Errors, Latency}")

# ================================================================ Cost engine
feature_slide(
    6, "FONCTIONNALITÉ 1/6", "Moteur de coût : qui dépense quoi",
    "Opacité des coûts : la facture LLM est un bloc unique. Impossible de dire quelle équipe ou "
    "application dépense le plus, ni de repérer qu'un modèle moins cher suffirait.",
    ["Coût calculé depuis les vrais tokens × prix déclaré par AIProvider.",
     "Ventilation par modèle, fournisseur, namespace, équipe, app.",
     "Coût moyen / requête et coût / token.",
     "Modèles vus sans pricing → comptés non pricés (data-quality).",
     "Moteur pur, sans dépendance K8s, ~95% de couverture de tests."],
    ["Chargeback / showback par centre de coût.",
     "Repère les modèles non tarifés (catalogue à compléter).",
     "Fondation des moteurs budget, break-even et du rapport.",
     "Pas de conversion multi-devise en MVP (devise dominante)."],
    formula="coût = (tokens / 1 000 000) × prixParMillion   —   pour l'input ET l'output, puis agrégé",
    metric="ai_finops_cost_eur · ai_finops_input/output_tokens · ai_finops_requests")

# ================================================================ 7 — Budget engine
feature_slide(
    7, "FONCTIONNALITÉ 2/6", "Moteur de budget : alerter avant le dépassement",
    "Sans garde-fou, on découvre le dépassement de budget IA après coup, sur la facture du mois. "
    "Aucun signal intermédiaire, aucune action graduée.",
    ["Compare la dépense d'une cible (ns/équipe/app) à un budget.",
     "usagePercent = round(spend / budget × 100).",
     "Phases : WithinBudget → Warning → Critical → Exceeded.",
     "Actions par seuil (onWarning / onCritical / onHardLimit).",
     "Modèle de repli (fallbackModelRef) : dégradation, pas blocage."],
    ["On voit venir le dépassement, pas après la facture.",
     "Dégradation gracieuse plutôt que coupure sèche.",
     "Métrique d'enforcement émise par action déclenchée.",
     "État lisible dans le .status (phase, % d'usage, dépense)."],
    formula="≥70% Warning   ·   ≥90% Critical   ·   ≥100% Exceeded   (seuils configurables par policy)",
    metric="ai_finops_budget_usage_percent{ns,policy} · ai_finops_enforcement_actions{policy,ns,app,mode,action} · Event K8s")

# ================================================================ 8 — Sovereignty engine
feature_slide(
    8, "FONCTIONNALITÉ 3/6", "Moteur de souveraineté : où part la donnée ?",
    "Souveraineté & conformité : en secteur régulé il faut savoir où sont traitées les données et "
    "si le fournisseur est autorisé — aujourd'hui géré au cas par cas, sans trace.",
    ["Confronte CHAQUE flux observé (ns/app → modèle → fournisseur).",
     "Zone interdite → critical ; hors zones autorisées → warning.",
     "Données sensibles vers fournisseur managé non habilité → warning.",
     "NormalizeZone : 'france'→FR, 'eastus'→US ; EU couvre les membres.",
     "Findings attribués au flux + nombre de requêtes concernées."],
    ["Détecte les flux non conformes, attribués à l'app fautive.",
     "Volume réel de requêtes à risque (pas juste un booléen).",
     "Traçabilité d'audit (RGPD / AI Act), pas une attestation.",
     "reportOnly = radar ; enforce = réécrit la route Envoy (slide 9)."],
    metric="ai_finops_sovereignty_findings · ai_finops_sovereignty_requests{ns,app,policy,sev} · ai_finops_cost_by_zone_eur{zone}")

# ================================================================ 9 — Enforcement (plan de données)
feature_slide(
    9, "PLAN DE CONTRÔLE → DONNÉES", "Enforcement : du constat à l'action dans la gateway",
    "Détecter ne suffit pas en secteur régulé : quand une app envoie ses données vers un fournisseur "
    "hors zone autorisée, il faut pouvoir RÉELLEMENT rerouter le trafic — pas seulement le signaler.",
    ["Modes gradués par policy : reportOnly → warn → enforce.",
     "Moteur d'enforcement pur (décisions) + actuateur de gateway.",
     "enforce réécrit l'AIGatewayRoute Envoy (client unstructured) :",
     "modèle interdit → backend conforme + bodyMutation du modèle.",
     "Réversible : backend d'origine sauvé en annotation, revert auto."],
    ["Reroute RÉEL validé : gpt-4o (US) → Mistral EU (Azure Foundry).",
     "Le fournisseur interdit n'est jamais atteint, la requête réussit.",
     "Budget : métrique d'enforcement émise par action déclenchée.",
     "Finalizers : revert propre à la suppression de la policy.",
     "Plan de données inchangé hors enforce (radar par défaut)."],
    metric="ai_finops_enforcement_actions{policy, namespace, application, mode, action, actuated}")

# ================================================================ 10 — Break-even engine
feature_slide(
    10, "FONCTIONNALITÉ 4/6", "Moteur break-even : API managée ou GPU auto-hébergé ?",
    "Mauvais arbitrage économique : les équipes ne savent pas à partir de quel volume héberger leur "
    "propre modèle sur GPU devient moins cher que l'API managée.",
    ["Coût managé/mois = coût tokens + frais fixes fournisseur.",
     "Coût auto-hébergé/mois = GPU + ops + stockage/réseau.",
     "Économie mensuelle, puis payback = migration / économie.",
     "Verdict : keep-managed / self-host (≤ seuil) / investigate.",
     "Extrapole l'usage observé à un mois (seuil payback déf. 6 mois)."],
    ["Décision d'investissement GPU fondée sur l'usage réel.",
     "Chiffre l'économie potentielle et le délai de retour.",
     "Coûts GPU/ops fournis par l'utilisateur (hypothèses explicites).",
     "Se recoupe avec la souveraineté (héberger en UE = 2 gains)."],
    formula="savings = managé − auto-hébergé   ·   payback (mois) = coût_migration / savings",
    metric="ai_finops_breakeven_savings_eur{namespace, analysis}")

# ================================================================ 10 — Recommandations
feature_slide(
    11, "FONCTIONNALITÉ 5/6", "Moteur de recommandation : passer du constat à l'action",
    "Mesurer ne suffit pas : il faut dire QUOI faire, chiffré et priorisé — sans noyer l'équipe "
    "sous des suggestions à faible gain.",
    ["Croise usage + catalogue + findings de souveraineté.",
     "cost-saving : trouve le modèle le moins cher pour chaque app.",
     "Filtre anti-bruit : suggéré seulement si gain ≥ 20% du coût.",
     "sovereignty (critical) : requêtes/coût vers provider non conforme.",
     "data-quality (warning) : modèle sans prix → créer AIModel/AIProvider."],
    ["Liste d'actions priorisée (sévérité, puis € d'économie).",
     "Chaque action attribuée à son app (namespace/application).",
     "Économie chiffrée en € sur la fenêtre observée.",
     "Action concrète : modèle courant → modèle recommandé.",
     "Moteur pur et testé (recommendationengine)."],
    metric="ai_finops_recommendations · _potential_savings_eur · _potential_savings_by_app_eur · _cost_saving_eur")

# ================================================================ 11 — Reporting
feature_slide(
    12, "FONCTIONNALITÉ 6/6", "Reporting : un dossier d'audit lisible et reproductible",
    "Les constats de coût et de souveraineté restent éparpillés. Il faut un livrable unique, "
    "lisible par un humain ET exploitable par l'outillage, pour préparer un audit.",
    ["Rend un rapport en Markdown (lisible) et JSON (contrat stable).",
     "Résumé exécutif : coût total, requêtes, tokens, coût/req, coût/token.",
     "Coût par modèle / fournisseur / équipe + top modèles.",
     "Findings de souveraineté + recommandations + limites & hypothèses.",
     "Publié dans un ConfigMap <nom>-report (GC automatique)."],
    ["Un seul livrable consolidé, daté et reproductible.",
     "Mentionne explicitement : préparation d'audit, pas attestation.",
     "Rappelle le mode reportOnly (aucune altération du trafic).",
     "Résultats aussi disponibles dans le .status de l'AIFinOpsReport."],
    metric="kubectl get cm <nom>-report -o jsonpath='{.data.report\\.md}'   (et report.json)")

# ================================================================ 13 — Shadow-AI (eBPF)
feature_slide(
    13, "PLAN INDÉPENDANT DE LA GATEWAY", "Shadow-AI : capter ce qui contourne la gateway",
    "Les collecteurs ne voient que le trafic qui PASSE par la gateway. Un pod qui appelle "
    "api.openai.com en direct échappe à tout — c'est le shadow-AI, le plus gros angle mort.",
    ["Tetragon (eBPF) observe l'egress par pod DANS le noyau.",
     "L'opérateur classe la destination par zone (EndpointToZone).",
     "Moteur pur shadowengine : zone interdite → critical.",
     "Indépendant de la gateway (tourne même sans AIGateway).",
     "Tetragon = DaemonSet standalone (tout CNI, pas Cilium)."],
    ["Ferme l'angle mort : trafic IA hors gouvernance détecté.",
     "Attribué au pod/namespace fautif, par zone de souveraineté.",
     "Données 100% réelles (connexions eBPF), jamais fabriquées.",
     "VALIDÉ LIVE sur AKS : finance → OpenAI US = critical.",
     "Backends futurs : Hubble, capture SNI, gRPC natif."],
    metric="ai_finops_shadow_ai_egress{namespace, application, zone, provider, severity}")

# ================================================================ 14 — Observabilité
s = slide()
title_bar(s, "TRANSVERSE", "Observabilité : 16 métriques ai_finops_* exposées", 14)
txt(s, Inches(0.55), Inches(1.3), Inches(12.2), Inches(0.4),
    [[("Vérifiées dans ", 13.5, False, DARKTX),
      ("internal/metrics/metrics.go", 13.5, True, GREEN),
      (" — exposées sur /metrics (:8080). Dashboard Grafana fourni.", 13.5, False, DARKTX)]])

def metric_panel(x, w, groups):
    y = Inches(1.85)
    for (head, col, lines) in groups:
        rect(s, x, y, w, Inches(0.42), col)
        txt(s, x+Inches(0.15), y, w-Inches(0.3), Inches(0.42),
            [[(head, 12.5, True, WHITE)]], anchor=MSO_ANCHOR.MIDDLE)
        y = Inches(y.inches + 0.42)
        for j, ln in enumerate(lines):
            bg = LIGHT if j%2==0 else WHITE
            rect(s, x, y, w, Inches(0.38), bg)
            txt(s, x+Inches(0.18), y, w-Inches(0.3), Inches(0.38),
                [[(ln, 11.5, False, DARKTX)]], anchor=MSO_ANCHOR.MIDDLE)
            y = Inches(y.inches + 0.38)
        y = Inches(y.inches + 0.12)

metric_panel(Inches(0.55), Inches(6.0), [
    ("Usage & coût  (label : namespace)", GREEN, [
        "ai_finops_requests",
        "ai_finops_input_tokens",
        "ai_finops_output_tokens",
        "ai_finops_cost_eur",
        "ai_finops_projected_monthly_cost_eur",
    ]),
    ("Souveraineté · enforcement · shadow-AI", BLUE, [
        "ai_finops_cost_by_zone_eur {zone}",
        "ai_finops_sovereignty_findings {ns,app,policy,sev}",
        "ai_finops_sovereignty_requests {ns,app,policy,sev}",
        "ai_finops_enforcement_actions {policy,ns,app,mode,action}",
        "ai_finops_shadow_ai_egress {ns,app,zone,provider,sev}",
    ]),
])
metric_panel(Inches(6.75), Inches(6.0), [
    ("Budget & optimisation", ORANGE, [
        "ai_finops_budget_usage_percent {ns,policy}",
        "ai_finops_potential_savings_eur",
        "ai_finops_potential_savings_by_app_eur {ns,app}",
        "ai_finops_cost_saving_eur {ns,app,current,reco}",
        "ai_finops_breakeven_savings_eur {ns,analysis}",
        "ai_finops_recommendations {type,ns,app,sev}",
    ]),
])
rrect(s, Inches(6.75), Inches(5.55), Inches(6.0), Inches(1.15), LIGHT)
txt(s, Inches(6.95), Inches(5.64), Inches(5.6), Inches(1.0),
    [[("À noter : ", 11.5, True, DARK),
      ("ai_finops_errors_total n'est PAS exposée — c'est une", 11.5, False, DARKTX)],
     [("métrique lue en entrée par le collecteur Prometheus,", 11.5, False, DARKTX)],
     [("pas une sortie de l'opérateur.", 11.5, False, DARKTX)]])
footer(s)

# ================================================================ 12 — Les 7 CRDs
s = slide()
title_bar(s, "API DÉCLARATIVE", "7 CRDs : la gouvernance comme objets Kubernetes (aiops.imperium.io/v1alpha1)", 15)
crds = [
    ("AIGateway", "aigw", "Gateway IA observée + mode de télémétrie + namespaces gouvernés."),
    ("AIProvider", "aiprov", "Fournisseur : région, dataResidency, managé ?, prix/M tokens, conformité."),
    ("AIModel", "—", "Modèle catalogué, lié à un provider ; tiers qualité/coût, données sensibles."),
    ("AIBudgetPolicy", "aibudget", "Budget €/période par cible + seuils + actions + modèle de repli."),
    ("AISovereigntyPolicy", "aisov", "Zones autorisées/interdites, données sensibles, audit, enforcementMode."),
    ("AIBreakEvenAnalysis", "aibreakeven", "Point mort managé vs auto-hébergé + recommandation."),
    ("AIFinOpsReport", "aireport", "Rapport consolidé généré (coûts, findings, recommandations)."),
]
y0 = Inches(1.45); rh = Inches(0.74)
for i,(name, short, desc) in enumerate(crds):
    y = y0 + i*rh
    bg = LIGHT if i%2==0 else WHITE
    rect(s, Inches(0.55), y, Inches(12.2), rh-Inches(0.07), bg)
    rect(s, Inches(0.55), y, Pt(4), rh-Inches(0.07), GREEN)
    txt(s, Inches(0.78), y, Inches(3.1), rh-Inches(0.07), [[(name, 14.5, True, DARK)]], anchor=MSO_ANCHOR.MIDDLE)
    txt(s, Inches(3.9), y, Inches(1.5), rh-Inches(0.07), [[(short, 12, True, BLUE)]], anchor=MSO_ANCHOR.MIDDLE)
    txt(s, Inches(5.45), y, Inches(7.2), rh-Inches(0.07), [[(desc, 12.5, False, DARKTX)]], anchor=MSO_ANCHOR.MIDDLE)
footer(s)

# ================================================================ 13 — Déploiement
s = slide()
title_bar(s, "MISE EN ŒUVRE", "Déploiement & exploitation", 16)
left = [
    ("Local (kind, sans build)", GREEN,
     ["make install        # CRDs", "make run            # manager", "kubectl apply -k config/samples/"]),
    ("Via Helm", BLUE,
     ["helm install greenops \\", "  charts/ai-sovereign-finops-operator \\", "  -n greenops-system --create-namespace"]),
    ("Démo réelle — trafic réel (1 commande)", ORANGE,
     ["cd automatisation/envoy-aigw && ./deploy.sh up", "# apps réelles → Envoy AI Gateway → LLM", "# tokens mesurés, pas de fake · dashboard :3000"]),
]
y = Inches(1.5)
for (t,c,lines) in left:
    rrect(s, Inches(0.55), y, Inches(6.1), Inches(1.55), WHITE, line=c)
    rect(s, Inches(0.55), y, Inches(6.1), Inches(0.42), c)
    txt(s, Inches(0.75), y+Inches(0.03), Inches(5.7), Inches(0.38), [[(t, 14, True, WHITE)]], anchor=MSO_ANCHOR.MIDDLE)
    txt(s, Inches(0.75), y+Inches(0.5), Inches(5.7), Inches(1.0),
        [[(l, 11, True, RGBColor(0x24,0x33,0x45))] for l in lines], space_after=4)
    y += Inches(1.72)
rrect(s, Inches(6.9), Inches(1.5), Inches(5.85), Inches(3.27), LIGHT)
txt(s, Inches(7.12), Inches(1.62), Inches(5.5), Inches(0.4), [[("Sécurisé par défaut", 15, True, DARK)]])
specs = ["Conteneur non-root, FS racine en lecture seule",
         "capabilities: drop ALL · seccomp RuntimeDefault",
         "RBAC au moindre privilège (généré)",
         "Secrets de gateway référencés, jamais copiés en status",
         "Service /metrics:8080 (--metrics-secure pour durcir)",
         "Chart : CRDs + Deployment + RBAC + ServiceMonitor opt."]
txt(s, Inches(7.12), Inches(2.1), Inches(5.5), Inches(2.6),
    [[("✓  " + l, 12.5, False, DARKTX)] for l in specs], space_after=8, ls=1.05)
rrect(s, Inches(0.55), Inches(5.4), Inches(12.2), Inches(1.3), DARK)
txt(s, Inches(0.8), Inches(5.52), Inches(11.7), Inches(1.1),
    [[("Stack 100% CNCF : ", 13.5, True, GREEN),
      ("Kubernetes · controller-runtime · Prometheus · Grafana · Helm · Kustomize · ArgoCD · Envoy", 13.5, False, WHITE)],
     [("Go 1.25 · Kubernetes ≥ 1.29 · Apache 2.0 · GitOps via ArgoCD (make up) ou chemin offline Helm (make local).", 13, False, MUTE)]])
footer(s)

# ================================================================ 14 — Synthèse
s = slide()
title_bar(s, "SYNTHÈSE", "Problème → fonctionnalité → impact", 17)
table = [
    ("Opacité des coûts", "Moteur de coût + catalogue défaut", "Chargeback dès l'install, sans CR"),
    ("Dépassement subi", "Moteur de budget (phases + actions)", "Alerte avant la facture"),
    ("Donnée hors contrôle", "Souveraineté + enforcement (reroute)", "Violation détectée ET reroutée (Envoy)"),
    ("Shadow-AI (hors gateway)", "Plan eBPF / Tetragon (shadowengine)", "Angle mort fermé (validé AKS)"),
    ("Arbitrage à l'aveugle", "Moteur break-even (payback)", "Décision GPU chiffrée"),
    ("Constats éparpillés", "Reporting (Markdown / JSON)", "Dossier d'audit reproductible"),
]
y0 = Inches(1.55); rh = Inches(0.86)
rect(s, Inches(0.55), y0, Inches(12.2), Inches(0.5), DARK)
hcols = [Inches(0.55), Inches(4.35), Inches(8.55)]; hwid = [Inches(3.8), Inches(4.2), Inches(4.2)]
for ci,h in enumerate(["Problème", "Fonctionnalité", "Impact"]):
    txt(s, hcols[ci]+Inches(0.15), y0, hwid[ci], Inches(0.5),
        [[(h, 14, True, GREEN if ci!=1 else BLUE)]], anchor=MSO_ANCHOR.MIDDLE)
for ri,(p,sln,imp) in enumerate(table):
    y = y0 + Inches(0.5) + ri*rh
    bg = LIGHT if ri%2==0 else WHITE
    rect(s, Inches(0.55), y, Inches(12.2), rh-Inches(0.04), bg)
    rect(s, Inches(0.55), y, Pt(4), rh-Inches(0.04), GREEN)
    txt(s, hcols[0]+Inches(0.15), y, hwid[0], rh, [[(p, 13, True, RED)]], anchor=MSO_ANCHOR.MIDDLE)
    txt(s, hcols[1]+Inches(0.15), y, hwid[1], rh, [[(sln, 12.5, False, DARKTX)]], anchor=MSO_ANCHOR.MIDDLE)
    txt(s, hcols[2]+Inches(0.15), y, hwid[2], rh, [[(imp, 12.5, True, RGBColor(0x12,0x7A,0x4D))]], anchor=MSO_ANCHOR.MIDDLE)
footer(s)

# ================================================================ 15 — Limites (honnêteté)
s = slide()
title_bar(s, "HONNÊTETÉ", "Limites & hypothèses (assumées)", 18)
lim = [
    ("Enforcement", "reportOnly (défaut) observe ; warn alerte ; enforce réécrit la route Envoy, réversible."),
    ("Pas d'attestation", "Prépare un dossier d'audit (RGPD/AI Act) ; ne garantit pas la conformité juridique."),
    ("Télémétrie", "prometheus · configmap · aigw (Envoy) réels ; fake opt-in explicite ; LiteLLM retiré."),
    ("Coûts déclaratifs", "Basés sur le pricing AIProvider ; modèles sans provider = non pricés."),
    ("Mono-devise", "Pas de conversion multi-devise ; devise dominante rapportée telle quelle."),
    ("Break-even simple", "Extrapolation linéaire ; coûts GPU/ops fournis par l'utilisateur ; payback déf. 6 mois."),
]
y0 = Inches(1.5); rh = Inches(0.84)
for i,(t,d) in enumerate(lim):
    y = y0 + i*rh
    rrect(s, Inches(0.55), y, Inches(12.2), rh-Inches(0.12), WHITE, line=RGBColor(0xD8,0xE0,0xE8))
    rect(s, Inches(0.55), y, Pt(5), rh-Inches(0.12), ORANGE)
    txt(s, Inches(0.8), y, Inches(3.0), rh-Inches(0.12), [[(t, 14, True, DARK)]], anchor=MSO_ANCHOR.MIDDLE)
    txt(s, Inches(3.9), y, Inches(8.7), rh-Inches(0.12), [[(d, 12.5, False, DARKTX)]], anchor=MSO_ANCHOR.MIDDLE)
footer(s)

# ================================================================ 16 — Clôture
s = slide()
rect(s, 0, 0, SW, SH, DARK)
rect(s, 0, 0, Inches(0.25), SH, GREEN)
txt(s, Inches(0.9), Inches(1.5), Inches(11.5), Inches(0.4),
    [[("EN UNE PHRASE", 14, True, GREEN)]])
txt(s, Inches(0.9), Inches(2.1), Inches(11.6), Inches(2.0),
    [[("Savoir, à tout instant et de façon prouvable,", 30, True, WHITE)],
     [("qui dépense quoi et où partent les données IA —", 30, True, WHITE)],
     [("sans fouiller le code de toutes les équipes.", 30, True, GREEN)]])
rect(s, Inches(0.9), Inches(4.6), Inches(6), Pt(3), GREEN)
txt(s, Inches(0.9), Inches(4.85), Inches(11.6), Inches(1.0),
    [[("Catalogue autonome  ·  2 plans souveraineté (gateway + eBPF)  ·  7 CRDs  ·  16 métriques + Grafana", 16, True, WHITE)],
     [("Validé bout-en-bout sur cluster réel (kind + AKS)  ·  enforce reroute Envoy  ·  shadow-AI eBPF  ·  100% CNCF  ·  Apache 2.0", 14, False, MUTE)]])
txt(s, Inches(0.9), Inches(6.1), Inches(11.6), Inches(0.5),
    [[("Prochaine étape : collecteur OTel générique · backend Tetragon gRPC + SNI · action block · multi-devise.", 14, False, MUTE)]])
txt(s, Inches(0.9), Inches(6.7), Inches(11.6), Inches(0.4),
    [[("Merci — questions ?", 16, True, GREEN)]])

out = "/mnt/c/Users/IhsenAlaya/Documents/ihsen/kubebuilder/greenops/AI-Sovereign-FinOps-Presentation.pptx"
prs.save(out)
print("OK ->", out, "| slides:", len(prs.slides._sldIdLst))
