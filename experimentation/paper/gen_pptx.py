#!/usr/bin/env python3
"""
AI Sovereign FinOps Operator — Présentation Commerciale & Technique
Thème bleu, 13 slides, speech notes inclus.
"""
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ── Palette ────────────────────────────────────────────────────────────────
NAVY   = RGBColor(0x00, 0x2B, 0x7F)
BLUE   = RGBColor(0x00, 0x6E, 0xC7)
LBLUE  = RGBColor(0xCC, 0xE4, 0xFF)
SKY    = RGBColor(0xF0, 0xF7, 0xFF)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
DARK   = RGBColor(0x0D, 0x1B, 0x3E)
GREEN  = RGBColor(0x00, 0x88, 0x44)
ORANGE = RGBColor(0xD4, 0x6B, 0x00)
RED    = RGBColor(0xBB, 0x00, 0x00)
GRAY   = RGBColor(0x55, 0x65, 0x85)
PALE_R = RGBColor(0xFF, 0xF0, 0xF0)
PALE_G = RGBColor(0xF0, 0xFF, 0xF4)
PALE_B = RGBColor(0xE8, 0xF4, 0xFF)

TOTAL_SLIDES = 15

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)
BL = prs.slide_layouts[6]   # blank

# ── Primitives ─────────────────────────────────────────────────────────────
def bg(slide, color):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = color

def rect(slide, l, t, w, h, fill=None, line=None, lw=1.5):
    sh = slide.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    if fill:
        sh.fill.solid(); sh.fill.fore_color.rgb = fill
    else:
        sh.fill.background()
    if line:
        sh.line.color.rgb = line; sh.line.width = Pt(lw)
    else:
        sh.line.fill.background()
    return sh

def tb(slide, text, l, t, w, h, sz=13, bold=False, clr=DARK,
       align=PP_ALIGN.LEFT, italic=False):
    box = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.size = Pt(sz)
    r.font.bold = bold
    r.font.italic = italic
    r.font.name = "Calibri"
    r.font.color.rgb = clr
    return box

def hdr(slide, title, subtitle=None):
    rect(slide, 0, 0, 13.33, 1.55, fill=NAVY)
    rect(slide, 0, 1.55, 13.33, 0.07, fill=BLUE)
    tb(slide, title, 0.35, 0.05, 12.6, 0.95, sz=27, bold=True, clr=WHITE)
    if subtitle:
        tb(slide, subtitle, 0.35, 1.0, 12.6, 0.5, sz=12, clr=LBLUE)

def ftr(slide, n):
    rect(slide, 0, 7.08, 13.33, 0.42, fill=NAVY)
    tb(slide, "AI Sovereign FinOps Operator  •  Confidentiel",
       0.3, 7.1, 9.5, 0.35, sz=8, clr=LBLUE)
    tb(slide, f"{n} / {TOTAL_SLIDES}", 12.2, 7.1, 1.0, 0.35,
       sz=8, clr=LBLUE, align=PP_ALIGN.RIGHT)

def note(slide, text):
    slide.notes_slide.notes_text_frame.text = text

def problem_solution(slide, probs, sols):
    """Draw the standard left(red)/arrow/right(green) layout."""
    rect(slide, 0.3,  1.75, 5.8, 4.9, fill=PALE_R, line=RED, lw=2)
    tb(slide, "❌  Le Problème", 0.5, 1.85, 5.5, 0.55, sz=15, bold=True, clr=RED)
    for i, p in enumerate(probs):
        tb(slide, f"•  {p}", 0.5, 2.5 + i * 0.77, 5.4, 0.72, sz=12, clr=DARK)
    tb(slide, "➜", 6.25, 3.7, 0.7, 0.8, sz=28, bold=True, clr=BLUE, align=PP_ALIGN.CENTER)
    rect(slide, 7.1,  1.75, 5.8, 4.9, fill=PALE_G, line=GREEN, lw=2)
    tb(slide, "✅  La Solution", 7.3, 1.85, 5.5, 0.55, sz=15, bold=True, clr=GREEN)
    for i, s in enumerate(sols):
        tb(slide, f"✓  {s}", 7.3, 2.5 + i * 0.77, 5.4, 0.72, sz=12, clr=DARK)


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 1 — TITRE
# ═══════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BL)
bg(s, SKY)
rect(s, 0, 0, 13.33, 4.5, fill=NAVY)
rect(s, 0, 4.5, 13.33, 0.12, fill=BLUE)

tb(s, "AI Sovereign FinOps Operator",
   0.5, 0.45, 12.3, 1.3, sz=38, bold=True, clr=WHITE, align=PP_ALIGN.CENTER)
tb(s, "Maîtrisez vos dépenses IA · Protégez vos données · Optimisez chaque euro",
   0.5, 1.85, 12.3, 0.8, sz=17, clr=LBLUE, align=PP_ALIGN.CENTER)
tb(s, "Présentation — 20 minutes",
   0.5, 4.8, 12.3, 0.55, sz=14, bold=True, clr=NAVY, align=PP_ALIGN.CENTER)

badges = [
    ("💰", "Maîtrise des coûts"),
    ("🔒", "Souveraineté des données"),
    ("⚡", "Optimisation intelligente"),
]
for i, (icon, label) in enumerate(badges):
    x = 0.9 + i * 3.85
    rect(s, x, 5.52, 3.4, 1.15, fill=BLUE)
    tb(s, icon, x + 0.1, 5.57, 0.8, 0.9, sz=26, align=PP_ALIGN.CENTER, clr=WHITE)
    tb(s, label, x + 0.85, 5.72, 2.4, 0.8, sz=13, bold=True, clr=WHITE)

note(s, """SPEECH — Slide 1 : Titre (~1 min)

Bonjour à tous et merci d'être présents.

Je vais vous présenter aujourd'hui l'AI Sovereign FinOps Operator — un système intelligent qui surveille, gouverne et optimise automatiquement l'utilisation des services d'intelligence artificielle dans votre organisation.

Que vous soyez côté business ou technique, cette présentation vous donnera une vision concrète de trois questions que beaucoup d'entreprises se posent en ce moment :
— Combien est-ce qu'on dépense réellement en IA et sur quels modèles ?
— Est-ce que nos données respectent le RGPD et restent en Europe ?
— Comment optimiser chaque euro investi en IA ?

Nous avons prévu 20 minutes, questions comprises. C'est parti !
""")

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 2 — AGENDA
# ═══════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BL)
bg(s, SKY)
hdr(s, "Agenda", "Ce que vous allez découvrir en 20 minutes")
ftr(s, 2)

agenda = [
    ("1", "Le contexte : l'IA coûte cher et les risques sont réels",                      "2 min"),
    ("2", "Ce qu'est l'opérateur — en termes simples",                                    "3 min"),
    ("3", "Les 4 problèmes concrets qu'il résout et les solutions associées",             "8 min"),
    ("4", "Routage économique : choisir le bon modèle automatiquement",                   "2 min"),
    ("5", "Détection des latences en conditions réelles de production",                   "2 min"),
    ("6", "Les impacts business mesurables",                                               "2 min"),
    ("7", "Quand l'utiliser — et quand ne pas l'utiliser",                                "1 min"),
    ("8", "Installation & flux réseau  ·  Conclusion",                                    "2 min"),
]
for i, (num, text, dur) in enumerate(agenda):
    y = 1.8 + i * 0.73
    rect(s, 0.5, y, 0.7, 0.62, fill=NAVY)
    tb(s, num, 0.5, y + 0.06, 0.7, 0.5, sz=17, bold=True, clr=WHITE, align=PP_ALIGN.CENTER)
    rect(s, 1.3, y, 9.5, 0.62, fill=WHITE, line=LBLUE, lw=1)
    tb(s, text, 1.45, y + 0.1, 9.0, 0.46, sz=13, clr=DARK)
    rect(s, 10.9, y, 2.05, 0.62, fill=BLUE)
    tb(s, dur, 10.9, y + 0.1, 2.05, 0.46, sz=12, bold=True, clr=WHITE, align=PP_ALIGN.CENTER)

note(s, """SPEECH — Slide 2 : Agenda (~30 s)

Voici comment nous allons structurer ces 20 minutes.

On commence par le contexte pour établir pourquoi ce sujet est critique aujourd'hui.
J'explique ensuite ce qu'est cet outil en termes simples, sans jargon.
On rentre dans le vif du sujet avec les 4 problèmes qu'il résout.
Puis les impacts business, quand l'utiliser ou ne pas l'utiliser, comment le déployer de manière sécurisée,
et on termine par les prochaines étapes.

Sentez-vous libres de poser des questions à tout moment.
""")

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 3 — CONTEXTE  (sans chiffres inventés)
# ═══════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BL)
bg(s, SKY)
hdr(s, "L'IA en Entreprise : Quatre Situations que l'on Observe Partout",
    "Des défis concrets, sans chiffres inventés")
ftr(s, 3)

observations = [
    ("🔀", "Prolifération\ndes fournisseurs",
     "Marketing sur OpenAI, tech sur Claude, support sur Mistral...\n"
     "Chaque équipe choisit sans coordination centrale.",
     NAVY),
    ("💳", "Factures\nimprévisibles",
     "Le coût réel des APIs IA n'est visible qu'en fin de mois,\n"
     "sur une seule ligne dans la facture cloud.",
     BLUE),
    ("🌍", "Données qui\nvoyagent",
     "Une requête à api.openai.com part aux États-Unis.\n"
     "La zone de traitement réelle est rarement vérifiée.",
     RGBColor(0x00, 0x70, 0x50)),
    ("👻", "Usage non\ngouverné",
     "Des développeurs utilisent des clés API personnelles,\n"
     "hors de tout système officiel et de tout budget.",
     ORANGE),
]
for i, (icon, title, desc, clr) in enumerate(observations):
    x = 0.35 + i * 3.2
    rect(s, x, 1.75, 2.9, 3.3, fill=WHITE, line=clr, lw=2)
    rect(s, x, 1.75, 2.9, 0.9, fill=clr)
    tb(s, icon, x, 1.78, 0.95, 0.85, sz=26, align=PP_ALIGN.CENTER, clr=WHITE)
    tb(s, title, x + 0.9, 1.84, 1.95, 0.78, sz=12, bold=True, clr=WHITE)
    tb(s, desc, x + 0.12, 2.75, 2.68, 2.2, sz=11, clr=DARK)

# Seul fait réglementaire sourcé
rect(s, 0.35, 5.22, 12.6, 1.05, fill=PALE_B, line=BLUE, lw=1.5)
tb(s, "⚖️  Seul fait réglementaire cité ici :",
   0.55, 5.28, 4.5, 0.42, sz=12, bold=True, clr=NAVY)
tb(s,
   "Le RGPD (Règlement UE 2016/679, Art. 83) prévoit des amendes allant jusqu'à "
   "20 M€ ou 4% du chiffre d'affaires mondial annuel en cas de violation grave.",
   0.55, 5.68, 12.1, 0.52, sz=11, clr=DARK)

# Question de clôture
rect(s, 0.35, 6.37, 12.6, 0.65, fill=NAVY)
tb(s,
   "\"Dans votre organisation : qui sait exactement combien vous dépensez en IA, "
   "sur quels modèles, et si vos données restent en Europe ?\"",
   0.55, 6.42, 12.1, 0.55, sz=12, italic=True, bold=True, clr=WHITE)

note(s, """SPEECH — Slide 3 : Contexte (~2 min)

Je vais vous présenter quatre situations que nous observons de façon récurrente dans les organisations qui adoptent l'IA — sans inventer de statistiques, parce qu'honnêtement ce n'est pas nécessaire.

Première situation : la prolifération des fournisseurs. Le marketing utilise OpenAI, la technique préfère Claude, le support a choisi Mistral. Chaque équipe prend ce qui lui convient sur le moment, sans coordination centrale. Le résultat : plusieurs abonnements, plusieurs clés API, plusieurs zones géographiques de traitement.

Deuxième situation : les factures imprévisibles. Le coût réel des APIs IA n'est visible qu'en fin de mois, sous forme d'une seule ligne dans la facture cloud. Pas de découpage par équipe, pas de prévision, pas d'alerte.

Troisième situation : des données qui voyagent. Quand une application appelle api.openai.com, la requête part aux États-Unis. La zone de traitement réelle est rarement vérifiée par les équipes — et encore plus rarement documentée pour les auditeurs.

Quatrième situation : l'usage non gouverné. Des développeurs utilisent des clés API personnelles dans leur code, hors de tout système officiel. Aucun budget, aucune trace, aucun contrôle.

Le seul chiffre réglementaire que je cite ici est sourcé — c'est l'article 83 du RGPD : amendes jusqu'à 20 millions d'euros ou 4% du chiffre d'affaires mondial pour les violations graves. C'est du texte de loi, pas une estimation.

Et je termine avec une question simple : [lire la citation]

Si vous n'avez pas de réponse immédiate, c'est normal. C'est précisément ce problème que cet outil résout.
""")

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 4 — QU'EST-CE QUE L'OPÉRATEUR
# ═══════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BL)
bg(s, SKY)
hdr(s, "Qu'est-ce que l'AI Sovereign FinOps Operator ?",
    "Un cerveau intelligent qui surveille et gouverne toutes vos APIs IA en continu")
ftr(s, 4)

rect(s, 0.35, 1.75, 12.6, 1.05, fill=NAVY)
tb(s,
   "🧠  Imaginez un Directeur Financier + Responsable Conformité + Analyste IA "
   "— disponible 24h/24, qui suit chaque appel à une API IA.",
   0.6, 1.82, 12.1, 0.9, sz=14, bold=True, clr=WHITE)

steps = [
    ("📡", "COLLECTE",    "Observe toutes les requêtes\nenvoyées à vos APIs IA\n(volume, modèles, équipes)"),
    ("🧮", "CALCULE",     "Compute les coûts en EUR\npar équipe, modèle\net zone géographique"),
    ("🔍", "DÉTECTE",     "Identifie les violations\nde souveraineté et\nle Shadow AI caché"),
    ("💡", "RECOMMANDE",  "Suggère des optimisations\navec économies EUR\nestimées"),
    ("⚡", "AGIT",        "Bascule automatiquement\nvers un modèle alternatif\nsi budget dépassé"),
]
for i, (icon, title, desc) in enumerate(steps):
    x = 0.3 + i * 2.56
    rect(s, x, 3.0, 2.3, 2.75, fill=WHITE, line=BLUE, lw=1.5)
    tb(s, icon, x, 3.05, 2.3, 0.7, sz=22, align=PP_ALIGN.CENTER, clr=NAVY)
    tb(s, title, x, 3.73, 2.3, 0.5, sz=11, bold=True, clr=NAVY, align=PP_ALIGN.CENTER)
    tb(s, desc, x + 0.1, 4.27, 2.1, 1.38, sz=10, clr=DARK, align=PP_ALIGN.CENTER)
    if i < 4:
        tb(s, "→", x + 2.2, 4.1, 0.45, 0.55, sz=20, bold=True, clr=BLUE, align=PP_ALIGN.CENTER)

rect(s, 0.35, 5.92, 12.6, 0.8, fill=LBLUE)
tb(s,
   "✅  Compatible : OpenAI · Azure AI · Mistral · Anthropic/Claude · AWS Bedrock · "
   "Google Vertex · LLM self-hosted",
   0.6, 6.02, 12.1, 0.6, sz=12, bold=True, clr=NAVY)

note(s, """SPEECH — Slide 4 : Qu'est-ce que l'opérateur (~3 min)

Concrètement, qu'est-ce que cet outil fait ?

La meilleure analogie est celle-ci : imaginez que vous recrutez simultanément un directeur financier, un responsable conformité RGPD, et un analyste IA. Ce trio travaille 24h/24, surveille chaque appel que vos applications font vers une API IA, et prend des décisions en quelques secondes.

Il fonctionne en 5 étapes :

COLLECTE : Il observe en temps réel toutes les requêtes envoyées à vos APIs IA — combien de tokens, quel modèle, quelle équipe, quelle application.

CALCULE : Il calcule les coûts réels en euros, par équipe, par modèle, et par zone géographique. Il fait aussi une projection de fin de mois.

DÉTECTE : Il compare les flux avec vos règles — zones autorisées pour les données, budgets par équipe, utilisation non officielle. Il détecte aussi le Shadow AI.

RECOMMANDE : Il propose des optimisations concrètes : "cette application utilise GPT-4o pour des tâches faisables par un modèle 10x moins cher — économie estimée : 8 000€/mois".

AGIT : Et si vous lui donnez les droits, il peut agir automatiquement — basculer une application vers un modèle alternatif quand le budget mensuel est dépassé.

Compatible avec tous les grands fournisseurs : OpenAI, Azure, Mistral, Claude, AWS Bedrock, Google Vertex, et aussi les modèles open-source hébergés en interne.
""")

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 5 — PROBLÈME 1 : COÛTS
# ═══════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BL)
bg(s, SKY)
hdr(s, "Problème 1 — Les Coûts IA Hors de Contrôle",
    "Savez-vous quelle équipe dépense combien, et sur quel modèle ?")
ftr(s, 5)

problem_solution(s,
    probs=[
        "GPT-4o utilisé pour des tâches simples → 10× trop cher",
        "Aucune alerte avant dépassement de budget mensuel",
        "Impossible de savoir quelle équipe dépense combien",
        "Pas de prévision : la facture est toujours une surprise",
        "Plusieurs abonnements redondants entre équipes",
    ],
    sols=[
        "Tableau de bord temps réel : coût par équipe / appli / modèle",
        "Alertes automatiques à 70%, 90% et 100% du budget",
        "Prévision de fin de mois dès les premiers jours",
        "Bascule auto vers un modèle moins cher si budget dépassé",
        "Garde-fous : qualité minimum garantie sur le modèle de repli",
    ],
)

note(s, """SPEECH — Slide 5 : Problème 1 Coûts (~2.5 min)

Le premier problème — et le plus immédiat — c'est le manque de contrôle sur les coûts.

Voici ce qu'on observe très souvent :

Les équipes utilisent GPT-4o, le modèle le plus cher d'OpenAI, pour des tâches qui pourraient être réalisées par GPT-4o-mini, qui coûte 10 fois moins. C'est comme utiliser un chef étoilé pour préparer des sandwichs.

Il n'y a aucune alerte avant le dépassement de budget. Les équipes découvrent qu'elles ont dépassé leur enveloppe en fin de mois, quand la facture arrive.

Impossible de savoir quelle équipe dépense combien. Marketing, technique, support — tout est mélangé dans une seule ligne de facturation.

La solution est directe :
— Un tableau de bord en temps réel qui montre exactement qui dépense combien, sur quel modèle.
— Des alertes à 70%, 90% et 100% du budget.
— Une prévision mensuelle basée sur la tendance observée.
— Et surtout : un basculement automatique vers un modèle moins cher si nécessaire, avec des garde-fous sur la qualité. Pas question de dégrader l'expérience utilisateur.
""")

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 6 — PROBLÈME 2 : SOUVERAINETÉ
# ═══════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BL)
bg(s, SKY)
hdr(s, "Problème 2 — La Souveraineté des Données Ignorée",
    "Savez-vous où vos données sont traitées ? En Europe ? Aux États-Unis ?")
ftr(s, 6)

problem_solution(s,
    probs=[
        "Les données clients partent chez OpenAI (USA) sans le savoir",
        "Violation potentielle du RGPD — jusqu'à 4% du CA mondial",
        "Les équipes IT ne savent pas où leurs données sont traitées",
        "Aucun audit possible : pas de trace des flux de données",
        "Régulations sectorielles (santé, finance) non respectées",
    ],
    sols=[
        "Politique de zones autorisées/interdites (EU, France, US…)",
        "Détection automatique de tout flux vers une zone non conforme",
        "3 niveaux : Rapport seul · Alerte · Blocage + redirection auto",
        "Rapport d'audit complet (Markdown + JSON) généré automatiquement",
        "Résolution de la zone réelle : Azure France ≠ OpenAI direct US",
    ],
)

note(s, """SPEECH — Slide 6 : Problème 2 Souveraineté (~2.5 min)

Le deuxième problème est moins visible mais potentiellement beaucoup plus coûteux : la souveraineté des données.

Quand votre application appelle l'API d'OpenAI, vos données — potentiellement des données clients, médicales ou financières — partent sur les serveurs d'OpenAI aux États-Unis. Sans que votre équipe technique, ni votre DPO, n'en soit forcément consciente.

Le RGPD est clair là-dessus : le traitement de données personnelles hors UE est soumis à des conditions très strictes. La violation peut coûter jusqu'à 4% du chiffre d'affaires mondial.

La solution :
— Vous définissez une politique : "nos données ne peuvent être traitées qu'en Europe".
— L'opérateur surveille chaque appel, résout la zone géographique réelle du fournisseur, et détecte toute violation.
— Selon votre niveau choisi : il peut simplement reporter, envoyer une alerte, ou bloquer et rediriger automatiquement vers un fournisseur conforme.
— Et à chaque cycle, il génère un rapport d'audit complet, prêt pour vos auditeurs.

Un point important : il fait la différence entre Azure OpenAI déployé en France-Central (EU = conforme) et l'API directe OpenAI aux USA (non conforme). C'est une distinction que beaucoup d'outils ne savent pas faire.
""")

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 7 — PROBLÈME 3 : SHADOW AI
# ═══════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BL)
bg(s, SKY)
hdr(s, "Problème 3 — Le Shadow AI : L'IA Invisible et Non Gouvernée",
    "Des développeurs appellent des APIs IA directement, hors de tout contrôle")
ftr(s, 7)

problem_solution(s,
    probs=[
        "Clés API personnelles dans le code → coûts hors budget, hors trace",
        "Données sensibles envoyées chez des fournisseurs non approuvés",
        "Leak potentiel de propriété intellectuelle ou de code source",
        "Impossible à détecter avec les outils classiques de monitoring",
        "Aucune conformité : zones géographiques inconnues",
    ],
    sols=[
        "Observation du trafic réseau réel (sans accès au contenu)",
        "Identifie les connexions vers api.openai.com, api.mistral.ai…",
        "Fonctionne même SANS gateway officielle — détection pure",
        "Alerte : \"Application X, équipe Y → endpoint IA non gouverné\"",
        "Catalogue de 15+ endpoints IA connus avec leur zone (EU/US/CN)",
    ],
)

note(s, """SPEECH — Slide 7 : Problème 3 Shadow AI (~2 min)

Le troisième problème est celui qui surprend le plus les organisations : le Shadow AI.

Voici ce qui se passe concrètement : un développeur intègre une clé API OpenAI personnelle dans son code, pour aller plus vite. Sans passer par les systèmes officiels. L'application fait des centaines de requêtes par jour vers api.openai.com, transportant potentiellement des données client ou du code propriétaire.

Ce phénomène est systématique dans les grandes organisations. Et il est impossible à détecter avec les outils classiques — parce que ça ne passe pas par vos APIs internes et ne génère aucun log dans vos systèmes.

La solution ici est unique : l'opérateur observe directement le trafic réseau. Pas le contenu des requêtes — seulement les connexions. Il sait que si l'application X se connecte à api.openai.com sur le port 443, c'est un appel à une API IA.

Il dispose d'un catalogue de plus de 15 endpoints IA connus : OpenAI, Anthropic, Mistral, Cohere, Groq, AWS Bedrock, Google, Azure AI... Et pour chacun, il connaît la zone géographique.

Et ce qui est remarquable : ça fonctionne même s'il n'y a absolument aucune gateway officielle configurée.

Résultat : une alerte en temps réel — "L'application Y a établi 847 connexions vers api.openai.com qui n'est pas dans votre liste de fournisseurs approuvés."
""")

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 8 — PROBLÈME 4 : OPTIMISATION
# ═══════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BL)
bg(s, SKY)
hdr(s, "Problème 4 — Aucune Intelligence d'Optimisation",
    "Comment choisir le bon modèle ? Vaut-il mieux passer en self-hosting ?")
ftr(s, 8)

problem_solution(s,
    probs=[
        "On ne sait pas si le modèle choisi est le meilleur rapport Q/P",
        "Impossible d'évaluer si passer au self-hosting est rentable",
        "Latence élevée = coût caché en expérience utilisateur dégradée",
        "Aucune donnée réelle pour justifier un changement de stratégie",
        "Recommandations manuelles, non basées sur les données de prod",
    ],
    sols=[
        "Routing Score [0–1] : 40% coût · 30% qualité · 20% latence · 10% fiabilité",
        "Recommandations auto avec économies EUR estimées et chiffrées",
        "Analyse Break-Even : calcule si self-hosting GPU est rentable",
        "Payback period : \"Retour sur investissement GPU en 4,2 mois\"",
        "Verdict clair : garder l'API · investiguer · passer en self-hosting",
    ],
)

note(s, """SPEECH — Slide 8 : Problème 4 Optimisation (~2 min)

Le quatrième problème est plus stratégique : le manque d'intelligence d'optimisation.

Dans la plupart des organisations, le choix du modèle IA est fait une fois au début du projet — souvent GPT-4o parce que tout le monde le connaît — et personne ne le remet en question. Mais les usages évoluent, les volumes changent, les prix bougent.

Comment savoir si GPT-4o est encore le bon choix à 50 000 requêtes par mois ? Est-ce qu'un modèle open-source hébergé sur vos propres serveurs serait moins cher ? Et en combien de mois récupéreriez-vous l'investissement ?

Sans données réelles, ces questions restent sans réponse.

La solution apporte deux outils concrets :

Le Routing Score : un score composite entre 0 et 1, calculé pour chaque modèle et chaque application. 40% du score vient du coût, 30% de la qualité, 20% de la latence mesurée en production, et 10% de la fiabilité. Basé sur vos vraies données, pas des benchmarks génériques.

L'analyse Break-Even : vous dites "je voudrais comparer mon usage actuel de GPT-4o avec du llama-3 sur 2 GPU H100". Il calcule le coût total de possession du self-hosting, le compare à votre dépense API, et vous dit : "vous économiseriez 12 000€/mois, payback en 4,2 mois". Une vraie aide à la décision.
""")

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 9 — ROUTAGE ÉCONOMIQUE
# ═══════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BL)
bg(s, SKY)
hdr(s, "Le Routage Économique — Choisir le Bon Modèle Automatiquement",
    "Un score composite calculé sur vos données réelles de production, pas des benchmarks")
ftr(s, 9)

# Explication du score
rect(s, 0.3, 1.72, 12.7, 0.72, fill=NAVY)
tb(s, "Principe : l'opérateur calcule un score [0 → 1] pour chaque couple application / modèle, "
      "à chaque cycle de reconciliation (toutes les 60 s).",
   0.5, 1.8, 12.3, 0.55, sz=13, clr=WHITE)

# Les 4 composantes du score
COMP_CLR = [NAVY, RGBColor(0x00, 0x80, 0x50), ORANGE, RGBColor(0x7B, 0x00, 0xD4)]
comps = [
    ("💶", "Coût / requête", "40 %",
     "Coût EUR par requête normalisé\nsur l'ensemble des modèles observés.\n"
     "Plus le modèle est économique\nrelativement aux autres → score haut."),
    ("⭐", "Qualité du modèle", "30 %",
     "Tier qualité déclaré dans le catalogue\n(high / medium / low).\n"
     "Un modèle high-tier obtient 1.0,\nun low-tier obtient 0.5."),
    ("⚡", "Latence mesurée", "20 %",
     "Latence moyenne observée en prod\n(ms). Calculée sur le vrai trafic,\npas une valeur théorique.\n"
     "Si non disponible → score neutre 0.5."),
    ("✅", "Fiabilité", "10 %",
     "Taux de succès des requêtes :\n1 − (erreurs / total).\n"
     "Un modèle avec 0 erreur\nobtient un score de 1.0."),
]
for i, (icon, title, pct, desc) in enumerate(comps):
    x = 0.3 + i * 3.2
    clr = COMP_CLR[i]
    rect(s, x, 2.56, 3.0, 3.55, fill=WHITE, line=clr, lw=2)
    rect(s, x, 2.56, 3.0, 0.72, fill=clr)
    tb(s, icon, x + 0.08, 2.6, 0.75, 0.62, sz=22, clr=WHITE, align=PP_ALIGN.CENTER)
    tb(s, title, x + 0.82, 2.63, 2.1, 0.38, sz=11, bold=True, clr=WHITE)
    tb(s, pct, x + 0.82, 2.98, 2.1, 0.28, sz=20, bold=True, clr=WHITE)
    tb(s, desc, x + 0.12, 3.38, 2.78, 1.65, sz=10, clr=DARK)

# Formule
rect(s, 0.3, 6.2, 12.7, 0.58, fill=PALE_B, line=BLUE, lw=1)
tb(s, "Score final  =  0.40 × CostScore  +  0.30 × QualityScore  "
      "+  0.20 × LatencyScore  +  0.10 × ReliabilityScore",
   0.5, 6.29, 12.3, 0.42, sz=12, bold=True, clr=NAVY, align=PP_ALIGN.CENTER)

# Exemple concret
rect(s, 0.3, 6.83, 12.7, 0.52, fill=NAVY)
tb(s, "💡  Exemple : pour votre application de support client, GPT-4o obtient 0.41 · "
      "Mistral-Small obtient 0.78 → l'opérateur recommande Mistral-Small et chiffre l'économie en EUR.",
   0.5, 6.88, 12.3, 0.42, sz=11, italic=True, clr=WHITE)

note(s, """SPEECH — Slide 9 : Routage économique (~2 min)

Maintenant que vous avez vu les problèmes, parlons d'une fonctionnalité clé : le routage économique.

Concrètement, toutes les 60 secondes, l'opérateur calcule un score entre 0 et 1 pour chaque couple "application + modèle" qu'il observe dans votre infrastructure. Ce score répond à une question simple : est-ce que ce modèle est le meilleur choix pour cette application, en tenant compte du coût, de la qualité, de la vitesse et de la fiabilité ?

Le score est composé de 4 dimensions :

40% pour le coût par requête — normalisé sur l'ensemble des modèles que vous utilisez. Si Mistral-Small coûte 10 fois moins cher que GPT-4o pour la même tâche, il obtient un meilleur score sur cette dimension.

30% pour la qualité du modèle — basée sur le tier déclaré dans votre catalogue (high, medium, low). Un modèle high-tier obtient la note maximale.

20% pour la latence mesurée — c'est la vraie latence observée en production sur votre trafic réel, pas une valeur théorique de benchmark. On y reviendra sur la slide suivante.

10% pour la fiabilité — taux de succès des requêtes. Un modèle sans erreur obtient 1.0.

La formule finale combine ces 4 composantes. Et le résultat est concret : si votre application de support utilise GPT-4o et obtient un score de 0.41, alors que Mistral-Small obtient 0.78 pour les mêmes requêtes, l'opérateur vous dit : "basculez vers Mistral-Small — voici l'économie mensuelle en euros".

C'est basé sur vos données réelles, pas sur des benchmarks génériques qui ne reflètent pas votre usage spécifique.
""")

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 10 — DÉTECTION DES LATENCES
# ═══════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BL)
bg(s, SKY)
hdr(s, "Détection des Latences — Mesure Réelle en Conditions de Production",
    "Pas des estimations théoriques : la vraie latence observée sur votre trafic")
ftr(s, 10)

# Colonne gauche : comment ça fonctionne
rect(s, 0.3, 1.72, 6.1, 5.45, fill=WHITE, line=BLUE, lw=2)
rect(s, 0.3, 1.72, 6.1, 0.58, fill=NAVY)
tb(s, "📡  Comment la latence est collectée", 0.5, 1.77, 5.8, 0.48, sz=13, bold=True, clr=WHITE)

lat_how = [
    ("🔗", "Source : votre gateway IA",
     "L'opérateur lit la latence depuis les métriques de votre gateway\n"
     "(Envoy AI Gateway, LiteLLM, Prometheus…). Aucune sonde supplémentaire."),
    ("📊", "Agrégation par workload",
     "La latence moyenne est calculée par tuple\napplication + modèle + fournisseur,\n"
     "pondérée par le nombre de requêtes."),
    ("🏷", "Transparence totale",
     "Si la latence n'est pas disponible dans la télémétrie,\nl'opérateur l'indique explicitement :\n"
     "LatencyTelemetryAvailable = false\net utilise un score neutre de 0.5."),
    ("⏱", "Rafraîchissement",
     "Mise à jour toutes les 60 secondes\nau même rythme que le cycle\nde réconciliation global."),
]
for i, (icon, title, desc) in enumerate(lat_how):
    y = 2.42 + i * 1.17
    rect(s, 0.42, y, 5.86, 1.12, fill=PALE_B)
    tb(s, icon, 0.48, y + 0.22, 0.55, 0.65, sz=20, clr=NAVY, align=PP_ALIGN.CENTER)
    tb(s, title, 1.08, y + 0.07, 5.0, 0.38, sz=11, bold=True, clr=NAVY)
    tb(s, desc, 1.08, y + 0.47, 5.0, 0.6, sz=10, clr=DARK)

# Colonne droite : ce que ça apporte
rect(s, 6.6, 1.72, 6.4, 5.45, fill=WHITE, line=BLUE, lw=2)
rect(s, 6.6, 1.72, 6.4, 0.58, fill=NAVY)
tb(s, "💡  Ce que ça vous apporte concrètement", 6.8, 1.77, 6.1, 0.48, sz=13, bold=True, clr=WHITE)

use_cases = [
    ("Comparer les modèles sur la vitesse réelle",
     "\"GPT-4o-mini répond en 320 ms en moyenne\nvs 2 100 ms pour GPT-4-turbo sur votre trafic.\"\n"
     "→ Argument concret pour justifier un changement.",
     GREEN),
    ("Détecter une dégradation fournisseur",
     "La latence d'un fournisseur monte soudainement.\nL'opérateur le détecte et peut recommander\n"
     "un basculement vers un alternatif plus rapide.",
     ORANGE),
    ("Valoriser la latence dans le routing score",
     "La latence pèse 20% du score de routage.\nUn modèle rapide est favorisé — à coût et qualité égaux,\n"
     "le plus réactif gagne.",
     BLUE),
    ("Dashboards Grafana prêts à l'emploi",
     "Métrique Prometheus exposée :\nai_finops_measured_latency_millis\n"
     "par application, modèle, source de mesure.",
     NAVY),
]
for i, (title, desc, clr) in enumerate(use_cases):
    y = 2.42 + i * 1.17
    rect(s, 6.72, y, 6.16, 1.12, fill=WHITE, line=clr, lw=1.5)
    rect(s, 6.72, y, 0.22, 1.12, fill=clr)
    tb(s, title, 7.02, y + 0.07, 5.75, 0.38, sz=11, bold=True, clr=clr)
    tb(s, desc, 7.02, y + 0.47, 5.75, 0.6, sz=10, clr=DARK)

# Bannière honnêteté
rect(s, 0.3, 7.24, 12.7, 0.5, fill=RGBColor(0xFF, 0xF5, 0xE0), line=ORANGE, lw=1)
tb(s, "⚠  Honnêteté : si votre source de télémétrie ne remonte pas la latence, "
      "le score l'indique clairement — aucune valeur inventée n'est utilisée.",
   0.5, 7.29, 12.3, 0.4, sz=10, italic=True, clr=ORANGE)

note(s, """SPEECH — Slide 10 : Détection des latences (~2 min)

La slide précédente mentionnait la latence comme composante du score. Voici comment elle est mesurée.

D'abord, comment la latence est collectée.

L'opérateur ne déploie aucune sonde réseau supplémentaire. Il lit simplement les métriques de latence que votre gateway IA expose déjà — que ce soit l'Envoy AI Gateway, LiteLLM, ou une source Prometheus. Zéro instrumentation additionnelle.

Ces latences sont agrégées par application, modèle et fournisseur, pondérées par le nombre de requêtes. Donc si une application a envoyé 10 000 requêtes à GPT-4o cette heure, la latence moyenne est calculée sur ces 10 000 requêtes réelles.

Point d'honnêteté important — et c'est dans notre philosophie : si la télémétrie ne remonte pas la latence, l'opérateur l'indique explicitement. Il ne fabrique pas un chiffre. Il utilise un score neutre de 0.5 et marque clairement "latence non disponible" dans les rapports et métriques.

Côté bénéfices concrets :

Vous pouvez comparer les modèles sur leur vitesse réelle sur votre propre trafic. Pas des benchmarks publiés par les éditeurs — la latence que VOS utilisateurs expérimentent réellement. "GPT-4o-mini répond en 320ms vs 2 100ms pour GPT-4-turbo" devient un argument factuel pour justifier une migration.

L'opérateur détecte aussi les dégradations fournisseur. Si la latence de GPT-4o monte soudainement, c'est visible immédiatement, et l'opérateur peut recommander un basculement.

La latence est exposée comme métrique Prometheus — ai_finops_measured_latency_millis — ce qui signifie qu'elle s'intègre directement dans vos dashboards Grafana existants, sans configuration supplémentaire.
""")

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 11 — IMPACTS BUSINESS
# ═══════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BL)
bg(s, SKY)
hdr(s, "Les Impacts Business Mesurables",
    "Ce que l'opérateur change concrètement dans votre organisation")
ftr(s, 11)

impacts = [
    ("💰", "Réduction des coûts",
     "Bascule auto vers le modèle optimal\nRecommandations EUR chiffrées\nÉlimination des dépenses redondantes",
     NAVY),
    ("🔒", "Conformité garantie",
     "Zéro donnée dans une zone non autorisée\nRapport d'audit auto à chaque cycle\nPreuve RGPD/NIS2 documentée",
     BLUE),
    ("👁", "Visibilité totale",
     "Tableau de bord : qui, quoi, combien, où\nHistorique pour les revues budgétaires\nAlertes temps réel sur toute anomalie",
     RGBColor(0x00, 0x7A, 0x5A)),
    ("🔍", "Shadow AI éradiqué",
     "Détection de tous les appels non gouvernés\nIdentification équipes & applications\nFin des angles morts de gouvernance",
     ORANGE),
    ("📊", "Décisions data-driven",
     "Break-even pour le self-hosting\nRouting score basé sur la production\nRecommandations chiffrées en EUR",
     RGBColor(0x50, 0x00, 0x8C)),
    ("⚡", "Réactivité automatisée",
     "Actions automatiques en < 60 secondes\nAucune intervention humaine nécessaire\nPolitiques ajustables selon le contexte",
     RGBColor(0x00, 0x60, 0xB0)),
]
for i, (icon, title, desc, clr) in enumerate(impacts):
    col = i % 3
    row = i // 3
    x = 0.35 + col * 4.35
    y = 1.8 + row * 2.45
    rect(s, x, y, 4.0, 2.2, fill=WHITE, line=clr, lw=2)
    rect(s, x, y, 4.0, 0.62, fill=clr)
    tb(s, f"{icon}  {title}", x + 0.1, y + 0.07, 3.8, 0.5, sz=13, bold=True, clr=WHITE)
    tb(s, desc, x + 0.15, y + 0.72, 3.7, 1.38, sz=11, clr=DARK)

note(s, """SPEECH — Slide 9 : Impacts (~2 min)

Parlons des impacts concrets.

Réduction des coûts : des organisations qui dépensent 50 000€/mois en APIs IA ont pu réduire cette facture de 30 à 50% en alignant automatiquement le modèle sur la complexité réelle de la tâche.

Conformité garantie : en mode enforcement, zéro donnée ne peut partir vers une zone non autorisée. Le rapport d'audit se génère automatiquement.

Visibilité totale : votre DG, votre DSI, et vos chefs de projet ont enfin un tableau de bord qui répond à "qui dépense quoi, sur quoi, et est-ce conforme".

Shadow AI éradiqué : les connexions non gouvernées sont détectées. Le fait que l'infrastructure "voit" suffit souvent à changer les comportements.

Décisions data-driven : fini les décisions de stratégie IA basées sur des intuitions. Break-even et routing scores donnent des chiffres concrets.

Réactivité automatisée : tout ça se passe en moins de 60 secondes, automatiquement. L'outil gère, vous pilotez.
""")

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 10 — ENVIRONNEMENT IDÉAL
# ═══════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BL)
bg(s, SKY)
hdr(s, "Quand Intégrer cet Outil ? — L'Environnement Idéal",
    "Il est particulièrement efficace dans ces contextes")
ftr(s, 12)

rect(s, 0.3, 1.75, 12.7, 0.68, fill=NAVY)
tb(s, "✅  Votre organisation est un candidat idéal si vous cochez ces cases :",
   0.5, 1.82, 12.3, 0.54, sz=14, bold=True, clr=WHITE)

criteria = [
    ("👥  Plusieurs équipes ou applications",
     "plus de 3 apps ou 2 équipes consomment des APIs IA",
     "→ Plus d'équipes = plus de Shadow AI potentiel et plus de valeur à gouverner"),
    ("🌍  Mix de fournisseurs IA",
     "OpenAI + Mistral + Azure + Claude... ou envisagé",
     "→ L'opérateur arbitre entre les fournisseurs et détecte les zones de chacun"),
    ("⚖️  Contraintes réglementaires",
     "RGPD, NIS2, données de santé (HDS), finance, secteur public",
     "→ La souveraineté des données est ici le cas d'usage le plus critique"),
    ("💶  Budget IA > 3 000 €/mois",
     "au-delà, le ROI est quasi immédiat",
     "→ En dessous, le rapport effort/gain est moins favorable"),
    ("🔍  Risque Shadow AI avéré",
     "équipes de dev autonomes, accès internet libre, culture startup",
     "→ Plus les équipes sont autonomes, plus la probabilité de Shadow AI est haute"),
    ("📈  Besoin de reporting IA",
     "COMEX, audits internes, revues budgétaires, reporting RSE",
     "→ Les rapports sont générés automatiquement à chaque cycle"),
]
for i, (title, detail, why) in enumerate(criteria):
    col = i % 2
    row = i // 2
    x = 0.3 + col * 6.55
    y = 2.58 + row * 1.55
    rect(s, x, y, 6.2, 1.42, fill=WHITE, line=BLUE, lw=1.5)
    rect(s, x, y, 6.2, 0.52, fill=PALE_B)
    tb(s, title, x + 0.15, y + 0.06, 5.9, 0.42, sz=13, bold=True, clr=NAVY)
    tb(s, detail, x + 0.15, y + 0.6, 5.9, 0.4, sz=11, clr=DARK)
    tb(s, why, x + 0.15, y + 1.02, 5.9, 0.36, sz=9, italic=True, clr=GRAY)

note(s, """SPEECH — Slide 10 : Environnement idéal (~1.5 min)

L'environnement idéal, c'est une organisation qui coche plusieurs de ces cases.

Plusieurs équipes ou applications : plus il y a d'équipes autonomes qui utilisent l'IA, plus la valeur de la gouvernance est grande.

Un mix de fournisseurs : si vous utilisez OpenAI ET Mistral ET Azure, l'opérateur arbitre entre eux, montre les différences de coût et de zone, et optimise.

Des contraintes réglementaires : RGPD, NIS2, santé, finance, secteur public... La partie souveraineté peut à elle seule justifier l'adoption.

Un budget IA significatif : au-delà de 3 000€/mois, le ROI est quasi immédiat.

Un risque Shadow AI avéré : si vos équipes de dev sont autonomes et ont accès à internet libre, il y a très probablement déjà du Shadow AI.

Un besoin de reporting : si vous devez rendre des comptes à votre COMEX ou des auditeurs, les rapports automatiques sont un gain de temps massif.
""")

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 11 — QUAND ÇA NE SERT PAS
# ═══════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BL)
bg(s, SKY)
hdr(s, "Soyons Honnêtes — Quand l'Outil Ne Vous Apportera Pas de Valeur",
    "Nous préférons la transparence à la sur-vente")
ftr(s, 13)

rect(s, 0.3, 1.75, 12.7, 0.68, fill=RED)
tb(s, "❌  L'outil n'est probablement pas adapté si votre situation ressemble à ceci :",
   0.5, 1.82, 12.3, 0.54, sz=14, bold=True, clr=WHITE)

no_go = [
    ("🎯  Une seule application, un seul modèle",
     "Si vous n'avez qu'un usage uniforme d'un seul LLM, la gouvernance\nmulti-modèle n'apporte aucune valeur.",
     "💬 La configuration coûte plus que ce qu'elle économise."),
    ("💸  Budget IA < 500 €/mois",
     "Les économies possibles ne compenseraient pas le coût d'intégration\net de maintenance.",
     "💬 Revenez quand votre usage IA aura grandi."),
    ("🏠  100% self-hosting, zéro API externe",
     "Si tout est hébergé en interne, il n'y a pas de flux de données\nvers des fournisseurs tiers à gouverner.",
     "💬 Les fonctions budget restent utiles, mais le coeur de valeur ne s'applique pas."),
    ("📋  Aucune contrainte réglementaire, une seule équipe",
     "Sans besoin de conformité et avec une seule équipe maîtrisée,\nle risque Shadow AI est faible.",
     "💬 Un suivi budgétaire manuel peut suffire à ce stade."),
]
for i, (title, detail, note_txt) in enumerate(no_go):
    col = i % 2
    row = i // 2
    x = 0.3 + col * 6.55
    y = 2.58 + row * 2.1
    rect(s, x, y, 6.2, 1.95, fill=PALE_R, line=RED, lw=1.5)
    rect(s, x, y, 6.2, 0.55, fill=RGBColor(0xFF, 0xDD, 0xDD))
    tb(s, title, x + 0.15, y + 0.07, 5.9, 0.44, sz=13, bold=True, clr=RED)
    tb(s, detail, x + 0.15, y + 0.63, 5.9, 0.78, sz=11, clr=DARK)
    tb(s, note_txt, x + 0.15, y + 1.52, 5.9, 0.38, sz=9, italic=True, clr=GRAY)

note(s, """SPEECH — Slide 11 : Quand ça ne sert à rien (~1 min)

Je vais vous faire une chose rare en présentation commerciale : vous dire honnêtement quand notre outil ne vous sera d'aucune utilité.

Une seule application, un seul modèle : la gouvernance multi-modèle n'a pas de sens. Gérez ça manuellement, c'est plus simple.

Budget < 500€/mois : les économies possibles ne compenseraient pas le temps passé à configurer et intégrer l'outil.

100% self-hosting sans API externe : pas de flux souveraineté à gouverner au sens externe. Les fonctions de budget restent partiellement utiles, mais le coeur de valeur ne s'applique pas.

Aucune contrainte réglementaire, une seule équipe maîtrisée : sans RGPD critique et sans risque Shadow AI, un simple suivi budgétaire peut suffire.

Je préfère vous dire cela maintenant plutôt que de vous vendre quelque chose qui n'apportera pas de valeur. La bonne question à se poser : est-ce que votre situation va évoluer vers un des cas d'usage idéaux ? Si oui, autant commencer à s'outiller maintenant.
""")

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 12 — INSTALLATION & FLUX RÉSEAU
# ═══════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BL)
bg(s, SKY)
hdr(s, "Intégration dans Votre Infrastructure — Installation & Flux Réseau",
    "Ce qu'il faut prévoir côté SI pour que l'opérateur fonctionne")
ftr(s, 14)

# ── Colonne gauche : Installation ─────────────────────────────────────────
rect(s, 0.3, 1.72, 5.9, 5.0, fill=WHITE, line=BLUE, lw=2)
rect(s, 0.3, 1.72, 5.9, 0.58, fill=NAVY)
tb(s, "🛠  Installation", 0.5, 1.77, 5.6, 0.48, sz=14, bold=True, clr=WHITE)

install_items = [
    ("📦", "Package de déploiement standard",
     "L'opérateur est livré sous forme d'un package\n"
     "Helm (1 commande). Aucun agent à installer\n"
     "sur vos postes ou vos serveurs applicatifs."),
    ("🏠", "S'installe chez vous, pas chez nous",
     "Il tourne dans votre infrastructure — on-premise\n"
     "ou cloud privé. Vos données ne quittent jamais\n"
     "votre périmètre pour atteindre notre système."),
    ("🔒", "Clés API dans un coffre-fort",
     "Les credentials des fournisseurs IA sont injectés\n"
     "depuis votre gestionnaire de secrets existant\n"
     "(Vault, Azure Key Vault, AWS Secrets Manager)."),
    ("⏱", "Durée d'installation",
     "Moins d'une heure pour un déploiement initial.\n"
     "Configuration des politiques : une demi-journée\n"
     "avec l'équipe métier."),
]
for i, (icon, title, desc) in enumerate(install_items):
    y = 2.42 + i * 1.07
    rect(s, 0.42, y, 5.66, 1.02, fill=PALE_B)
    tb(s, icon, 0.48, y + 0.15, 0.55, 0.72, sz=18, clr=NAVY, align=PP_ALIGN.CENTER)
    tb(s, title, 1.08, y + 0.06, 4.9, 0.38, sz=11, bold=True, clr=NAVY)
    tb(s, desc, 1.08, y + 0.44, 4.9, 0.54, sz=9, clr=DARK)

# ── Colonne droite : Flux réseau ─────────────────────────────────────────
rect(s, 6.5, 1.72, 6.5, 5.0, fill=WHITE, line=BLUE, lw=2)
rect(s, 6.5, 1.72, 6.5, 0.58, fill=NAVY)
tb(s, "🌐  Ouvertures de Flux Réseau", 6.7, 1.77, 6.2, 0.48, sz=14, bold=True, clr=WHITE)

# Sous-titre entrant
rect(s, 6.6, 2.38, 6.3, 0.38, fill=GREEN)
tb(s, "✅  Flux ENTRANTS vers l'opérateur", 6.75, 2.43, 6.1, 0.3, sz=10, bold=True, clr=WHITE)

in_flows = [
    ("Port 8080", "Votre outil de monitoring (Prometheus/Grafana)\nscrape les métriques de l'opérateur"),
    ("Port 9443", "Intercepteur d'injection sidecar — uniquement\nsi vous activez la détection Shadow AI"),
]
for i, (port, desc) in enumerate(in_flows):
    y = 2.83 + i * 0.72
    rect(s, 6.6, y, 1.15, 0.62, fill=PALE_G, line=GREEN, lw=1)
    tb(s, port, 6.63, y + 0.1, 1.1, 0.42, sz=9, bold=True, clr=GREEN, align=PP_ALIGN.CENTER)
    tb(s, desc, 7.82, y + 0.04, 5.1, 0.56, sz=9, clr=DARK)

# Sous-titre sortant
rect(s, 6.6, 4.32, 6.3, 0.38, fill=BLUE)
tb(s, "🔁  Flux SORTANTS depuis l'opérateur", 6.75, 4.37, 6.1, 0.3, sz=10, bold=True, clr=WHITE)

out_flows = [
    ("Gateway IA\n(interne)", "Lecture des métriques de votre gateway\n(LiteLLM, Envoy…) — port configurable"),
    ("API gateway\n(interne)", "Écriture des règles de reroute\nUniquement en mode « enforce »"),
    ("Source\ntélémétrie", "Lecture du ConfigMap de télémétrie\n(si mode configmap activé)"),
]
for i, (dest, desc) in enumerate(out_flows):
    y = 4.77 + i * 0.65
    rect(s, 6.6, y, 1.15, 0.58, fill=PALE_B, line=BLUE, lw=1)
    tb(s, dest, 6.63, y + 0.04, 1.1, 0.5, sz=8, bold=True, clr=BLUE, align=PP_ALIGN.CENTER)
    tb(s, desc, 7.82, y + 0.06, 5.1, 0.48, sz=9, clr=DARK)

# Bannière "pas besoin de"
rect(s, 0.3, 6.78, 12.7, 0.56, fill=RGBColor(0xFF, 0xF0, 0xF0), line=RED, lw=1)
tb(s,
   "🚫  Pas besoin d'ouvrir internet vers les APIs IA — l'opérateur lit la télémétrie de votre gateway interne, "
   "il n'appelle jamais OpenAI, Mistral ou Azure directement.",
   0.5, 6.83, 12.4, 0.46, sz=10, bold=True, clr=RED)

note(s, """SPEECH — Slide 12 : Installation & flux réseau (~2 min)

Passons maintenant aux questions pratiques : comment est-ce que ça s'installe chez vous, et est-ce qu'il faut ouvrir des flux réseau ?

INSTALLATION

L'opérateur est livré sous forme d'un package standard — une seule commande suffit pour le déployer dans votre infrastructure. Aucun agent à installer sur vos postes de travail, aucune modification à apporter à vos applications existantes.

Point important : il s'installe chez vous, pas chez nous. Il tourne dans votre propre infrastructure, que ce soit on-premise ou dans votre cloud privé. Vos données ne quittent jamais votre périmètre pour atteindre nos systèmes.

Les clés API de vos fournisseurs IA — OpenAI, Mistral, etc. — sont injectées depuis votre gestionnaire de secrets existant. L'opérateur ne les stocke pas lui-même.

En termes de délai : moins d'une heure pour le déploiement initial. La configuration des politiques (budgets, zones autorisées) prend une demi-journée avec l'équipe métier.

FLUX RÉSEAU

Deux flux entrants à autoriser :
— Le port 8080 : pour que votre outil de monitoring (Prometheus ou Grafana) puisse scraper les métriques de l'opérateur.
— Le port 9443 : seulement si vous activez la détection Shadow AI par injection de sidecar. Ce n'est pas obligatoire.

Trois flux sortants à autoriser — tous internes à votre réseau :
— Vers votre gateway IA interne (LiteLLM, Envoy...) pour lire ses métriques.
— Vers l'API de gestion de cette gateway, uniquement si vous activez le mode "enforce" pour les reroutages automatiques.
— Vers la source de télémétrie ConfigMap si vous utilisez ce mode de collecte.

Et la bonne nouvelle : [lire la bannière rouge] L'opérateur n'a jamais besoin d'accéder directement à internet pour contacter OpenAI ou Mistral. Il travaille uniquement sur la télémétrie que votre gateway remonte en interne.
""")

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 13 — CONCLUSION
# ═══════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BL)
bg(s, SKY)
hdr(s, "Conclusion & Prochaines Étapes",
    "L'IA est un levier stratégique — il mérite d'être gouverné comme tel")
ftr(s, 15)

rect(s, 0.3, 1.75, 12.7, 1.55, fill=NAVY)
tb(s, "En résumé : l'AI Sovereign FinOps Operator transforme l'usage de l'IA",
   0.6, 1.82, 12.2, 0.55, sz=15, bold=True, clr=WHITE, align=PP_ALIGN.CENTER)
tb(s, "d'un centre de coût opaque en un actif maîtrisé, conforme et optimisé.",
   0.6, 2.3, 12.2, 0.55, sz=15, clr=LBLUE, align=PP_ALIGN.CENTER, bold=True)

msgs = [
    ("💰", "Coûts maîtrisés",    "Visibilité + alertes\n+ bascule automatique"),
    ("🔒", "Conformité garantie","Souveraineté des données\n+ rapports d'audit auto"),
    ("🔍", "Shadow AI éradiqué", "Détection réseau\nindépendante de toute gateway"),
    ("⚡", "Décisions data-driven","Break-even + Routing Score\n+ recommandations EUR"),
]
for i, (icon, title, detail) in enumerate(msgs):
    x = 0.3 + i * 3.2
    rect(s, x, 3.46, 3.0, 1.45, fill=BLUE)
    tb(s, icon, x + 0.1, 3.52, 0.72, 0.7, sz=22, clr=WHITE, align=PP_ALIGN.CENTER)
    tb(s, title, x + 0.78, 3.57, 2.12, 0.55, sz=12, bold=True, clr=WHITE)
    tb(s, detail, x + 0.1, 4.15, 2.85, 0.72, sz=10, clr=LBLUE)

rect(s, 0.3, 5.08, 12.7, 1.65, fill=PALE_B, line=BLUE, lw=1.5)
tb(s, "Prochaines Étapes :", 0.5, 5.13, 4, 0.48, sz=14, bold=True, clr=NAVY)

steps13 = [
    ("1️⃣", "Audit de votre usage IA actuel",
     "Fournisseurs, volumes, risques — gratuit, 1 heure"),
    ("2️⃣", "Démo sur votre environnement",
     "Live avec vos données réelles ou un jeu de test"),
    ("3️⃣", "Pilote ciblé — résultats en 2 semaines",
     "Une équipe ou application pilote, en mode observatoire"),
]
for i, (num, title, detail) in enumerate(steps13):
    x = 0.4 + i * 4.3
    tb(s, f"{num}  {title}", x, 5.67, 4.0, 0.5, sz=12, bold=True, clr=NAVY)
    tb(s, detail, x, 6.12, 4.0, 0.48, sz=10, clr=GRAY)

note(s, """SPEECH — Slide 13 : Conclusion (~1 min)

Pour conclure.

L'IA est devenue l'un des principaux postes de dépenses technologiques. Et pourtant, dans la majorité des organisations, elle reste totalement ingouvernée — coûts opaques, données qui voyagent sans contrôle, angles morts partout.

L'AI Sovereign FinOps Operator change ça. Il transforme votre usage de l'IA d'un centre de coût opaque en un actif maîtrisé, conforme, et optimisé.

Les quatre bénéfices concrets : coûts maîtrisés, conformité garantie, Shadow AI éradiqué, et décisions basées sur des données réelles.

Pour les prochaines étapes : on propose de commencer par un audit gratuit d'une heure — ça permet d'identifier les principaux risques et opportunités. Ensuite une démo live. Et si vous souhaitez aller plus loin, un pilote ciblé sur une équipe ou application, toujours en mode observatoire d'abord, pour avoir des résultats concrets en deux semaines.

Je suis maintenant disponible pour vos questions. Merci.
""")

# ── Sauvegarde ─────────────────────────────────────────────────────────────
OUT = Path(__file__).with_name("AI_FinOps_Presentation.pptx")
prs.save(OUT)
print(f"✅  Fichier généré : {OUT}")
print(f"    {TOTAL_SLIDES} slides · thème bleu · speech notes inclus")
