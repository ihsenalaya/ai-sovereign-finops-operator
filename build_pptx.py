#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""AI Sovereign FinOps Operator — Présentation 20 min
   Slides : bullets courts (3-7 mots)
   Notes  : speech naturel du présentateur
   Esprit : problème → solution → impact
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from lxml import etree

DARK   = RGBColor(0x0F, 0x1B, 0x2D)
GREEN  = RGBColor(0x1E, 0xC8, 0x7A)
BLUE   = RGBColor(0x3B, 0x82, 0xF6)
RED    = RGBColor(0xE5, 0x3E, 0x3E)
ORANGE = RGBColor(0xF5, 0x9E, 0x0B)
GREY   = RGBColor(0x5B, 0x66, 0x70)
LIGHT  = RGBColor(0xF4, 0xF7, 0xFA)
REDBG  = RGBColor(0xFD, 0xEC, 0xEC)
GRNBG  = RGBColor(0xE9, 0xF8, 0xF1)
BLUBG  = RGBColor(0xEA, 0xF1, 0xFB)
ORGBG  = RGBColor(0xFE, 0xF3, 0xE2)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
DARKTX = RGBColor(0x1A, 0x24, 0x30)
MUTE   = RGBColor(0xB8, 0xC6, 0xD4)

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height
BLANK  = prs.slide_layouts[6]

# ── helpers ──────────────────────────────────────────────────────────────
def new_slide(): return prs.slides.add_slide(BLANK)

def rect(s, x, y, w, h, color, line=None):
    shp = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    shp.fill.solid(); shp.fill.fore_color.rgb = color
    if line is None: shp.line.fill.background()
    else: shp.line.color.rgb = line; shp.line.width = Pt(1)
    shp.shadow.inherit = False; return shp

def rrect(s, x, y, w, h, color, line=None, lw=1.25):
    shp = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    shp.fill.solid(); shp.fill.fore_color.rgb = color
    if line is None: shp.line.fill.background()
    else: shp.line.color.rgb = line; shp.line.width = Pt(lw)
    shp.shadow.inherit = False; return shp

def txt(s, x, y, w, h, runs, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, sa=4, ls=1.0):
    tb = s.shapes.add_textbox(x, y, w, h); tf = tb.text_frame
    tf.word_wrap = True; tf.vertical_anchor = anchor
    tf.margin_left = Pt(3); tf.margin_right = Pt(3)
    tf.margin_top = Pt(1); tf.margin_bottom = Pt(1)
    for i, para in enumerate(runs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align; p.space_after = Pt(sa); p.space_before = Pt(0)
        p.line_spacing = ls
        for (t, sz, b, c, *rest) in para:
            r = p.add_run(); r.text = t
            r.font.size = Pt(sz); r.font.bold = b
            r.font.color.rgb = c; r.font.italic = rest[0] if rest else False
            r.font.name = "Calibri"
    return tb

def set_notes(slide, speech):
    ns = slide.notes_slide
    tf = ns.notes_text_frame
    tf.text = speech

def title_bar(s, kicker, title, num):
    rect(s, 0, 0, SW, Inches(1.12), DARK)
    rect(s, 0, Inches(1.12), SW, Pt(4), GREEN)
    txt(s, Inches(0.55), Inches(0.14), Inches(10.5), Inches(0.3),
        [[(kicker, 12, True, GREEN)]])
    txt(s, Inches(0.55), Inches(0.42), Inches(11.2), Inches(0.62),
        [[(title, 24, True, WHITE)]])
    c = s.shapes.add_shape(MSO_SHAPE.OVAL, SW-Inches(1.0), Inches(0.31), Inches(0.5), Inches(0.5))
    c.fill.solid(); c.fill.fore_color.rgb = GREEN
    c.line.fill.background(); c.shadow.inherit = False
    c.text_frame.text = str(num)
    pr = c.text_frame.paragraphs[0]; pr.alignment = PP_ALIGN.CENTER
    pr.runs[0].font.size = Pt(16); pr.runs[0].font.bold = True
    pr.runs[0].font.color.rgb = DARK

def footer(s, t="AI Sovereign FinOps Operator  ·  v0.5.2  ·  Apache 2.0"):
    txt(s, Inches(0.55), SH-Inches(0.38), Inches(12), Inches(0.3),
        [[(t, 9, False, GREY)]])

def bullets(s, x, y, w, h, items, color=DARKTX, size=13.5, ls=1.25):
    """Render a list of short bullet strings."""
    runs = [[(f"▸  {it}", size, False, color)] for it in items]
    txt(s, x, y, w, h, runs, sa=6, ls=ls)

def feature_slide(num, kicker, title,
                  prob_head, prob_bullets,
                  sol_bullets,
                  impact_bullets,
                  metric=None,
                  speech=""):
    s = new_slide()
    title_bar(s, kicker, title, num)

    # ── PROBLÈME ──────────────────────────────────────────────────────────
    top = Inches(1.38)
    rrect(s, Inches(0.55), top, Inches(12.2), Inches(1.22), REDBG, line=RED, lw=1.0)
    txt(s, Inches(0.78), top+Inches(0.08), Inches(2.5), Inches(0.3),
        [[("⚠  PROBLÈME", 11, True, RED)]])
    txt(s, Inches(0.78), top+Inches(0.38), Inches(11.7), Inches(0.75),
        [[(prob_head, 15, True, DARKTX)]])

    # ── DEUX COLONNES ─────────────────────────────────────────────────────
    cy = Inches(2.75); ch = Inches(3.1)
    rrect(s, Inches(0.55), cy, Inches(6.0), ch, GRNBG, line=GREEN, lw=1.0)
    txt(s, Inches(0.78), cy+Inches(0.1), Inches(5.5), Inches(0.32),
        [[("✅  CE QUE FAIT L'OPÉRATEUR", 11, True, GREEN)]])
    bullets(s, Inches(0.75), cy+Inches(0.48), Inches(5.6), Inches(2.5), sol_bullets, DARKTX)

    rrect(s, Inches(6.75), cy, Inches(6.0), ch, BLUBG, line=BLUE, lw=1.0)
    txt(s, Inches(6.98), cy+Inches(0.1), Inches(5.5), Inches(0.32),
        [[("💡  IMPACT BUSINESS", 11, True, BLUE)]])
    bullets(s, Inches(6.95), cy+Inches(0.48), Inches(5.6), Inches(2.5), impact_bullets, DARKTX)

    # ── MÉTRIQUE ──────────────────────────────────────────────────────────
    if metric:
        by = Inches(6.02)
        rrect(s, Inches(0.55), by, Inches(12.2), Inches(0.62), DARK)
        txt(s, Inches(0.78), by+Inches(0.1), Inches(11.8), Inches(0.42),
            [[("Métrique : ", 11, True, GREEN), (metric, 11, False, MUTE)]])

    footer(s)
    set_notes(s, speech)
    return s

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 1 — COVER
# ═══════════════════════════════════════════════════════════════════════════
s = new_slide()
rect(s, 0, 0, SW, SH, DARK)
rect(s, 0, 0, Inches(0.28), SH, GREEN)
rect(s, 0, SH-Pt(4), SW, Pt(4), GREEN)
txt(s, Inches(0.9), Inches(1.5), Inches(11.5), Inches(0.45),
    [[("KUBERNETES OPERATOR  ·  GOUVERNANCE IA", 13, True, GREEN)]])
txt(s, Inches(0.9), Inches(2.1), Inches(11.5), Inches(1.5),
    [[("AI Sovereign", 56, True, WHITE)],
     [("FinOps Operator", 56, True, GREEN)]])
txt(s, Inches(0.9), Inches(4.05), Inches(11.6), Inches(0.5),
    [[("Coût · Souveraineté · Budget · Routage · Shadow AI — en temps réel", 16, False, MUTE)]])
txt(s, Inches(0.9), Inches(5.1), Inches(11.6), Inches(0.5),
    [[("Commerciaux & Développeurs  ·  20 min", 12, False, GREY)]])
txt(s, Inches(0.9), Inches(6.4), Inches(11.6), Inches(0.5),
    [[("11 CRDs  ·  23 métriques  ·  3 providers réels  ·  Grafana radar 5D  ·  v0.5.2  ·  Apache 2.0", 13, True, GREEN)]])
set_notes(s, """Bonjour à tous. Je vais vous présenter en 20 minutes l'AI Sovereign FinOps Operator — un opérateur Kubernetes open source que nous avons développé pour répondre à un problème très concret : quand les équipes adoptent des LLM en production, personne ne sait combien ça coûte, si les données restent en France, et qui appelle quoi en dehors du gateway officiel.

La présentation s'adresse aux deux publics : les commerciaux et les décideurs, avec les problèmes business et les impacts ; et les développeurs, avec l'architecture technique et les CRDs.

Commençons par le problème.""")

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 2 — LES 3 ANGLES MORTS
# ═══════════════════════════════════════════════════════════════════════════
s = new_slide()
title_bar(s, "LE PROBLÈME", "3 angles morts quand l'IA arrive plus vite que la gouvernance", 2)
cols = [
    (RED,    "💸  Coût incontrôlable",
     ["Facture globale — qui consomme quoi ?",
      "Aucune décomposition par équipe",
      "Le CFO découvre le dépassement trop tard"]),
    (ORANGE, "🔒  Conformité aveugle",
     ["RGPD & AI Act : données en zone EU/FR",
      "Aucun outil ne valide la zone du provider",
      "Découverte à l'audit — après incident"]),
    (RED,    "👻  Shadow AI invisible",
     ["Équipes appellent api.openai.com direct",
      "Zéro visibilité · zéro budget · zéro conformité",
      "DSI découvre ça sur la facture"]),
]
cw = Inches(3.95); gap = Inches(0.2)
for i, (col, head, buls) in enumerate(cols):
    x = Inches(0.55) + i*(cw+gap)
    rect(s, x, Inches(1.42), cw, Inches(0.52), col)
    txt(s, x+Inches(0.15), Inches(1.48), cw-Inches(0.3), Inches(0.4),
        [[(head, 14, True, WHITE)]], align=PP_ALIGN.CENTER)
    rrect(s, x, Inches(1.98), cw, Inches(3.5), LIGHT, line=col, lw=0.8)
    bullets(s, x+Inches(0.15), Inches(2.12), cw-Inches(0.25), Inches(3.2),
            buls, DARKTX, 13)
rrect(s, Inches(0.55), Inches(5.65), Inches(12.2), Inches(1.05), GRNBG, line=GREEN)
txt(s, Inches(0.78), Inches(5.75), Inches(11.8), Inches(0.88),
    [[("→  La solution : ", 14, True, GREEN),
      ("un opérateur Kubernetes qui transforme ces décisions ad hoc", 14, False, DARKTX)],
     [("    en politiques versionnées, auditables et réconciliées en continu.", 14, False, DARKTX)]],
    ls=1.1)
footer(s)
set_notes(s, """Ces trois problèmes coexistent dans la quasi-totalité des entreprises qui utilisent des LLM en production.

Premier angle mort : le coût. La facture Azure ou OpenAI arrive globale. Impossible de savoir quelle équipe a consommé quoi. Le CFO reçoit une facture de 50 000 euros sans pouvoir imputer.

Deuxième angle mort : la conformité. Le RGPD et l'AI Act imposent que les données restent en zone européenne. Mais personne ne valide que le modèle utilisé en production respecte cette contrainte. On découvre ça lors d'un audit — après l'incident.

Troisième angle mort : le Shadow AI. Des équipes contournent le gateway officiel et appellent api.openai.com directement. Zéro visibilité, zéro budget, zéro conformité. Le DSI l'apprend sur la facture.

Notre réponse : un opérateur Kubernetes qui gère tout ça de manière déclarative, auditée, et réconciliée en continu.""")

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 3 — ARCHITECTURE
# ═══════════════════════════════════════════════════════════════════════════
s = new_slide()
title_bar(s, "COMMENT", "Plan de données (gateway) + Plan de contrôle (opérateur)", 3)
rrect(s, Inches(0.55), Inches(1.42), Inches(12.2), Inches(1.45), BLUBG)
txt(s, Inches(0.75), Inches(1.5), Inches(10), Inches(0.3),
    [[("PLAN DE DONNÉES — Envoy AI Gateway", 11, True, BLUE)]])
bx = Inches(0.92)
for i, (label, col) in enumerate([
    ("Apps\nrh · finance · legal · marketing", BLUE),
    ("Envoy AI Gateway\ngpt-france-mini · gpt-us-mini\nmistral-large-latest", BLUE),
    ("Providers\nAzure OpenAI FR  (France ✅)\nAzure OpenAI US  (US ❌)\nMistral EU        (EU ✅)", GREY),
]):
    w = Inches(3.6); x = bx + i*(w+Inches(0.45))
    rrect(s, x, Inches(1.9), w, Inches(0.88), WHITE, line=col)
    txt(s, x+Inches(0.12), Inches(1.95), w-Inches(0.2), Inches(0.8),
        [[(label, 11, False, DARKTX)]], ls=1.05)
    if i < 2:
        txt(s, x+w, Inches(2.18), Inches(0.45), Inches(0.38),
            [[("→", 18, True, col)]], align=PP_ALIGN.CENTER)
txt(s, Inches(3.0), Inches(2.93), Inches(8), Inches(0.3),
    [[("gen_ai_client_token_usage  ·  gen_ai_server_request_duration  (OTel → Prometheus)", 10, True, BLUE)]],
    align=PP_ALIGN.CENTER)
rrect(s, Inches(0.55), Inches(3.35), Inches(12.2), Inches(2.12), DARK)
txt(s, Inches(0.75), Inches(3.42), Inches(10), Inches(0.3),
    [[("PLAN DE CONTRÔLE — Opérateur (11 CRDs · 23 métriques)", 11, True, GREEN)]])
for i, eng in enumerate(["cost engine", "budget engine", "sovereignty engine",
                          "routing score", "shadow AI engine", "reco engine"]):
    w = Inches(1.88); x = Inches(0.82)+i*(w+Inches(0.08))
    rrect(s, x, Inches(3.85), w, Inches(0.82), RGBColor(0x1A,0x2E,0x4A), line=GREEN, lw=0.8)
    txt(s, x+Inches(0.08), Inches(3.88), w-Inches(0.1), Inches(0.75),
        [[(eng, 11, True, GREEN)]], align=PP_ALIGN.CENTER)
txt(s, Inches(0.75), Inches(4.85), Inches(11.8), Inches(0.38),
    [[(".status CRDs  ·  23 métriques ai_finops_*  ·  Grafana radar 5D", 11, False, MUTE)]])
rrect(s, Inches(0.55), Inches(5.38), Inches(12.2), Inches(1.28), GRNBG, line=GREEN)
bullets(s, Inches(0.72), Inches(5.5), Inches(12.0), Inches(1.1), [
    "Tetragon / eBPF — détecte les appels LLM hors gateway (Shadow AI) au niveau noyau",
    "Reconciliation toutes les 30s — .status mis à jour en temps réel sur chaque CRD",
], DARKTX, 12.5)
footer(s)
set_notes(s, """L'architecture est en deux plans bien séparés.

Le plan de données, c'est Envoy AI Gateway. Il reçoit le trafic des applications, route vers le bon provider, et émet des métriques OTel — tokens réels, durée de requête — par namespace et par application, grâce à un sidecar injecté automatiquement.

Le plan de contrôle, c'est notre opérateur. Il lit ces métriques, applique les politiques déclarées dans les CRDs, calcule les scores, et écrit les résultats dans le .status de chaque objet Kubernetes. Prometheus scrappe ces métriques, Grafana les affiche.

Et en parallèle, Tetragon tourne en DaemonSet au niveau noyau — il observe les connexions TCP sur le port 443 et détecte les appels LLM qui contournent le gateway. C'est le Shadow AI engine.""")

# ═══════════════════════════════════════════════════════════════════════════
# SLIDES 4-11 — FONCTIONNALITÉS
# ═══════════════════════════════════════════════════════════════════════════

feature_slide(4, "FONCTIONNALITÉ 1", "Attribution du coût LLM par équipe et application",
    prob_head="Facture globale — impossible d'imputer par équipe",
    prob_bullets=[],
    sol_bullets=[
        "Sidecar injecte x-greenops-namespace automatiquement",
        "Gateway tag les métriques par namespace / app",
        "Tokens réels × prix liste Azure → coût EUR",
        "AIFinOpsReport : coût observé + projeté",
    ],
    impact_bullets=[
        "Chargeback exact par équipe dès J+1",
        "Identification des top consommateurs",
        "Base chiffrée pour arbitrage Make vs Buy",
        "Rapport généré automatiquement",
    ],
    metric="ai_finops_cost_eur{namespace, application, model}  ·  ai_finops_projected_monthly_cost_eur",
    speech="""Premier problème : la facture arrive globale. Impossible de savoir quelle équipe RH, finance ou legal consomme le plus.

Notre solution : un sidecar est injecté automatiquement dans chaque namespace. Il ajoute les en-têtes x-greenops-namespace et x-greenops-app à chaque requête. Le gateway Envoy les lit et tag les métriques OTel par namespace et par application.

L'opérateur multiplie les tokens réels mesurés par les prix officiels Azure, et écrit le coût en euros dans le .status de l'AIFinOpsReport.

Impact concret : vous avez le chargeback exact par équipe dès le lendemain. Vous savez que l'équipe finance a consommé 42 euros ce mois-ci, contre 8 euros pour RH. Et vous avez une projection mensuelle basée sur le run-rate observé.""")

feature_slide(5, "FONCTIONNALITÉ 2", "Budget mensuel avec enforcement automatique",
    prob_head="Budget IA dépassé sans alerte, sans frein",
    prob_bullets=[],
    sol_bullets=[
        "AIBudgetPolicy : budgetEUR + seuils warning/critical",
        "Projection mensuelle en run-rate continu",
        "Phases : WithinBudget → Warning → Critical → Exceeded",
        "enforce : blockOrRequireApproval",
    ],
    impact_bullets=[
        "Alerte avant le dépassement, pas après",
        "Blocage automatique ou approbation humaine",
        "Demo live : finance → phase Critical",
        "Audit trail : qui a dépassé, quand, combien",
    ],
    metric="ai_finops_budget_usage_percent{namespace, application}  ·  ai_finops_projected_monthly_cost_eur",
    speech="""Deuxième problème : un budget fixé en janvier est dépassé en mars sans que personne ne soit alerté.

On déclare un AIBudgetPolicy avec le budget mensuel en euros et trois seuils : warning à 70%, critical à 90%, et hardLimit à 100%. L'opérateur projette la dépense mensuelle en continu — il extrapole le run-rate sur 30 jours.

Quand le seuil est franchi, l'opérateur change la phase de l'objet Kubernetes : WithinBudget, Warning, Critical, Exceeded. En mode enforce, il peut bloquer les appels ou exiger une approbation humaine.

Dans la démo live, l'équipe finance a un budget serré — vous verrez la phase passer Critical en quelques minutes de trafic réel.""")

feature_slide(6, "FONCTIONNALITÉ 3", "Contrôle de souveraineté : FR ✅  EU ✅  US ❌",
    prob_head="RGPD & AI Act non vérifiés en production",
    prob_bullets=[],
    sol_bullets=[
        "AISovereigntyPolicy : allowedZones=[FR, EU]",
        "Chaque AIProvider porte sa zone réelle",
        "3 zones en demo : France · EU · US",
        "Modes : reportOnly → warn → enforce",
    ],
    impact_bullets=[
        "Azure OpenAI France (FR) → ✅ conforme",
        "Mistral EU Foundry (EU) → ✅ conforme",
        "Azure OpenAI US (US) → ❌ violation critique",
        "Findings exportés · audit trail complet",
    ],
    metric="ai_finops_sovereignty_findings{namespace, application, severity}  ·  ai_finops_cost_by_zone_eur{zone}",
    speech="""Troisième problème : le RGPD et l'AI Act imposent que les données restent en zone européenne. Mais rien ne valide ça en production.

On déclare une AISovereigntyPolicy avec les zones autorisées — FR et EU — et les zones interdites — US et CN. Chaque AIProvider dans le catalogue porte sa zone réelle : dataResidency fr, eu, ou us.

Dans la démo, vous voyez trois zones simultanément : Azure OpenAI France Central — conforme. Mistral EU sur Azure Foundry — conforme. Azure OpenAI US East — violation critique, severity=critical.

L'opérateur émet des métriques avec le namespace, l'application, et la sévérité de la violation. En mode enforce, il peut bloquer le trafic vers les zones interdites.""")

feature_slide(7, "FONCTIONNALITÉ 4", "Routing Score 5D — radar Grafana en temps réel",
    prob_head="Choix du modèle LLM : intuitif, non outillé",
    prob_bullets=[],
    sol_bullets=[
        "5 scores 0→1 : coût · qualité · latence · fiabilité · souveraineté",
        "RoutingScore = 0.40×coût + 0.30×qualité + 0.20×latence + 0.10×fiabilité",
        "Souveraineté = gate dur : score=0 si zone interdite",
        "Radar Grafana 5D mis à jour en continu",
    ],
    impact_bullets=[
        "Décision fondée sur 5 critères simultanément",
        "Visualisation radar : profil instantané par modèle",
        "Recommandation auto : Mistral EU > OpenAI US",
        "Aucun modèle non-souverain ne peut scorer >0",
    ],
    metric="ai_finops_routing_score · ai_finops_cost_score · ai_finops_quality_score · ai_finops_latency_score · ai_finops_sovereignty_score",
    speech="""Quatrième problème : choisir un modèle LLM se fait souvent sur le prix ou la réputation. Personne ne synthétise les 5 dimensions simultanément.

L'opérateur calcule en continu 5 scores normalisés entre 0 et 1 pour chaque modèle : le score de coût inversé — moins cher, plus élevé — la qualité basée sur le tier déclaré, la latence mesurée par le gateway, la fiabilité sur les dernières requêtes, et la souveraineté.

Ce dernier score est un gate dur : si la zone est interdite, le routing score total est 0, peu importe les autres dimensions.

Le résultat s'affiche dans le radar Grafana en temps réel — vous voyez instantanément que Mistral EU a un meilleur profil souveraineté que OpenAI US, et l'opérateur génère la recommandation automatiquement.""")

feature_slide(8, "FONCTIONNALITÉ 5", "Shadow AI — détection eBPF sans code change",
    prob_head="Appels LLM directs hors gateway — invisibles",
    prob_bullets=[],
    sol_bullets=[
        "Tetragon DaemonSet : observe tcp_connect:443 au niveau noyau",
        "Hôtes LLM connus : api.openai.com · anthropic.com · etc.",
        "Connexion capturée → shadow_ai_egress mis à jour",
        "Opérateur classe severity=critical (zone US hors gateway)",
    ],
    impact_bullets=[
        "Détection < 30 secondes sans changer le code",
        "Demo live : finance/shadow-ai-rogue détecté",
        "Panel Grafana : quelle app, quel host, combien",
        "Pas besoin de modifier les applications",
    ],
    metric="ai_finops_shadow_ai_egress{namespace, application, host}",
    speech="""Cinquième problème : des équipes appellent api.openai.com directement, sans passer par le gateway. Zéro visibilité, zéro conformité, zéro budget.

Notre solution utilise Tetragon, un DaemonSet eBPF qui observe les connexions TCP au niveau noyau Linux. Quand un pod ouvre une connexion sur le port 443 vers un host LLM connu — api.openai.com, anthropic.com, et autres — Tetragon capture l'événement.

L'opérateur reçoit cet événement, identifie le namespace et l'application, et met à jour la métrique shadow_ai_egress avec severity=critical.

Dans la démo, vous voyez finance/shadow-ai-rogue être détecté en moins de 30 secondes — sans aucune modification du code de l'application.""")

feature_slide(9, "FONCTIONNALITÉ 6", "AIChangeRequest — approbation humaine avant rerouting",
    prob_head="Opérateur qui agit seul en prod : risque perçu",
    prob_bullets=[],
    sol_bullets=[
        "AIChangeRequest : demande déclarative de changement",
        "Cycle : Pending → Approved → Actuated",
        "Actuation uniquement sur approbation explicite",
        "Rejected / Expired si délai dépassé",
    ],
    impact_bullets=[
        "Circuit d'approbation SOC 2 / ISO 27001",
        "Audit trail : qui · quoi · quand · pourquoi",
        "GitOps : la demande est versionnée dans git",
        "Compatible secteurs régulés : banque · santé",
    ],
    metric="kubectl get aichangerequest  |  .status.phase: Pending|Approved|Actuated|Rejected",
    speech="""Sixième problème : quand on parle d'un opérateur qui reroute automatiquement le trafic de production, les équipes ops ont peur. Et c'est légitime.

L'AIChangeRequest est notre réponse : toute action sur le routage doit passer par une demande déclarative. Elle démarre en phase Pending, un opérateur humain l'approuve — ou la rejette — et seulement à ce moment-là l'actuation se déclenche.

Si personne n'approuve dans le délai configuré, elle expire automatiquement. Toutes ces transitions sont horodatées et auditables.

C'est la fonctionnalité qui permet d'utiliser l'opérateur dans des environnements régulés — banque, santé, défense — avec un circuit de validation conforme SOC 2 ou ISO 27001.""")

feature_slide(10, "FONCTIONNALITÉ 7", "Routage déclaratif : canary 80/20 et reroute immédiat",
    prob_head="Modifier le routage LLM = intervention manuelle risquée",
    prob_bullets=[],
    sol_bullets=[
        "AIRoutingPolicy : règles déclaratives avec poids",
        "Canary : 80% gpt-france-mini / 20% mistral-large",
        "AIRouteOverride : reroute immédiat et temporaire",
        "Révocation : kubectl delete → retour automatique",
    ],
    impact_bullets=[
        "Zero-downtime routing — sans redémarrage gateway",
        "Rollback en 10 secondes",
        "A/B testing modèles sans toucher au code",
        "TTL configurable : override auto-expirant",
    ],
    metric="kubectl get airoutingpolicy · kubectl get airouteoverride  |  .status.active",
    speech="""Septième problème : pour basculer du trafic entre deux modèles LLM, il faut aujourd'hui modifier des configurations gateway à la main — sans rollback propre, sans historique.

L'AIRoutingPolicy permet de déclarer des règles de routage avec des poids. Par exemple : 80% du trafic vers gpt-france-mini, 20% vers mistral-large-latest. Le gateway reçoit ces règles sans redémarrage.

L'AIRouteOverride est pour les situations d'urgence : un reroute immédiat et temporaire. Quand on le supprime, le trafic revient automatiquement à la route normale. Rollback en 10 secondes.

Ces deux CRDs permettent de faire du A/B testing entre modèles sans toucher au code des applications.""")

feature_slide(11, "FONCTIONNALITÉ 8", "AIQualityGate — bloquer les modèles sous seuil",
    prob_head="Modèle moins cher = dégradation silencieuse de qualité",
    prob_bullets=[],
    sol_bullets=[
        "AIQualityGate : qualityTier minimum requis",
        "Tiers : low · medium · high · premium",
        "Violation si modèle sélectionné < seuil",
        "Intégré au routing score (dimension quality)",
    ],
    impact_bullets=[
        "Garantie contractuelle de qualité par app",
        "Prévention des substitutions silencieuses",
        "Visible dans le radar : colonne quality_score",
        "Alerte avant plainte utilisateur final",
    ],
    metric="ai_finops_quality_score{namespace, application, model}",
    speech="""Huitième problème : quand on optimise les coûts, on est tenté de basculer vers un modèle moins cher. Mais si la qualité chute, c'est l'utilisateur final qui le découvre en premier.

L'AIQualityGate permet de déclarer un seuil minimal de qualité par application. Si le modèle sélectionné est en dessous de ce tier, l'opérateur génère une violation et l'indique dans le .status.

Ce score de qualité est intégré comme l'une des 5 dimensions du routing score — avec un poids de 30%. Un modèle bon marché mais de mauvaise qualité sera pénalisé dans le score global.

Le panel Grafana montre la colonne quality_score dans le radar, et vous voyez immédiatement quel modèle est sous-performant.""")

feature_slide(12, "FONCTIONNALITÉ 9", "Break-even & recommandations automatiques",
    prob_head="Économies potentielles invisibles — pas d'outil de comparaison",
    prob_bullets=[],
    sol_bullets=[
        "AIBreakEvenAnalysis : actuel vs alternatif",
        "Calcul : (prix_actuel − prix_alternatif) × tokens_réels",
        "Recommandation auto si savings > seuil",
        "Exporté dans .status.monthlySavingsEUR",
    ],
    impact_bullets=[
        "Quick wins chiffrés par équipe",
        "Justification ROI pour basculer vers Mistral EU",
        "Recommandation : rh peut économiser X€/mois",
        "Décision Make vs Buy sur données réelles",
    ],
    metric="ai_finops_potential_savings_eur{namespace, application}  ·  ai_finops_cost_saving_eur",
    speech="""Neuvième problème : les équipes ne savent pas qu'un modèle dix fois moins cher existe pour leur cas d'usage. Les économies restent invisibles.

L'AIBreakEvenAnalysis compare le modèle actuellement utilisé avec un alternatif, sur les tokens réellement observés — pas des estimations. Si le modèle alternatif est moins cher, l'opérateur calcule les économies mensuelles et les écrit dans le .status.

Par exemple : si l'équipe RH passe de gpt-france-mini à mistral-large-latest, l'opérateur calcule que ça représente X euros d'économies par mois, basé sur les tokens réels du mois dernier.

Cette recommandation est automatiquement exportée en métrique Prometheus et visible dans le dashboard Grafana.""")

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 13 — DÉMO LIVE
# ═══════════════════════════════════════════════════════════════════════════
s = new_slide()
title_bar(s, "DÉMO LIVE", "Ce que vous voyez dans Grafana en ce moment", 13)
rows = [
    (GREEN, "Azure OpenAI France Central", "rh · legal",          "FR ✅ conforme",   "gpt-france-mini"),
    (BLUE,  "Mistral EU — Azure Foundry",  "marketing",           "EU ✅ conforme",   "mistral-large-latest"),
    (RED,   "Azure OpenAI US East",         "finance",             "US ❌ violation",  "gpt-us-mini"),
    (GREY,  "Shadow AI — eBPF",             "finance/shadow-rogue","US ❌ critique",   "api.openai.com direct"),
]
hdrs = ["Provider réel", "Applications", "Zone", "Modèle"]
hx   = [Inches(0.55), Inches(3.62), Inches(6.85), Inches(9.38)]
hw   = [Inches(2.9),  Inches(3.05), Inches(2.35), Inches(3.5)]
for i, (h, x, w) in enumerate(zip(hdrs, hx, hw)):
    rect(s, x, Inches(1.42), w-Inches(0.08), Inches(0.36), DARK)
    txt(s, x+Inches(0.1), Inches(1.47), w-Inches(0.15), Inches(0.28),
        [[(h, 11, True, GREEN)]])
for r, (col, prov, apps, zone, model) in enumerate(rows):
    y = Inches(1.82) + r*Inches(0.9)
    bg = GRNBG if col==GREEN else (BLUBG if col==BLUE else (REDBG if col==RED else LIGHT))
    rrect(s, Inches(0.55), y, Inches(12.28), Inches(0.84), bg, line=col, lw=0.8)
    rect(s, Inches(0.55), y, Pt(5), Inches(0.84), col)
    for val, x, w in zip([prov, apps, zone, model], hx, hw):
        txt(s, x+Inches(0.12), y+Inches(0.2), w-Inches(0.2), Inches(0.5),
            [[(val, 12.5, val==prov, DARKTX)]])
rrect(s, Inches(0.55), Inches(5.55), Inches(12.2), Inches(1.12), DARK)
txt(s, Inches(0.78), Inches(5.65), Inches(11.8), Inches(0.95),
    [[("Grafana → http://localhost:3000", 14, True, GREEN)],
     [("Radar 5D  ·  Budget phases  ·  Sovereignty findings  ·  Shadow AI bargauge", 12, False, MUTE)]],
    ls=1.2)
footer(s)
set_notes(s, """Voici ce que vous voyez en direct dans Grafana.

Quatre flux de trafic réels, trois providers, trois zones différentes.

Azure OpenAI France Central — gpt-france-mini — utilisé par les équipes RH et legal. Zone France, conforme. Score de souveraineté maximum.

Mistral EU sur Azure Foundry — mistral-large-latest — utilisé par marketing. Zone EU, conforme. Le modèle européen par excellence.

Azure OpenAI US East — gpt-us-mini — utilisé par finance. Zone US, violation critique. Vous voyez le findings counter augmenter en temps réel dans Grafana.

Et le Shadow AI : finance/shadow-ai-rogue appelle api.openai.com directement, sans passer par le gateway. Tetragon l'a détecté via eBPF — c'est la barre rouge dans le panel Shadow AI.

Ouvrons Grafana.""")

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 14 — 11 CRDs
# ═══════════════════════════════════════════════════════════════════════════
s = new_slide()
title_bar(s, "API DÉCLARATIVE", "11 CRDs — la gouvernance IA comme objets Kubernetes natifs", 14)
crds = [
    (GREEN,  "CATALOGUE",      [("AIGateway",          "endpoint Envoy + telemetry mode"),
                                 ("AIProvider",          "type · zone · prix · conformité"),
                                 ("AIModel",             "modelName · qualityTier · servesNamespace")]),
    (BLUE,   "COÛT & BUDGET",  [("AIFinOpsReport",      "coût observé · projeté · findings"),
                                 ("AIBudgetPolicy",      "budgetEUR · seuils · enforcement"),
                                 ("AIBreakEvenAnalysis", "actuel vs alternatif · savings")]),
    (ORANGE, "ROUTAGE",        [("AIRoutingPolicy",     "règles · poids · TTL"),
                                 ("AIRouteOverride",     "reroute immédiat · auto-révocation"),
                                 ("AIChangeRequest",     "Pending→Approved→Actuated")]),
    (RED,    "QUALITÉ & SÉC",  [("AIQualityGate",       "qualityTier minimum requis"),
                                 ("AISovereigntyPolicy", "zones autorisées/interdites · mode")]),
]
cw = Inches(2.95); gap = Inches(0.1)
for i, (col, cat, items) in enumerate(crds):
    x = Inches(0.55)+i*(cw+gap)
    rect(s, x, Inches(1.42), cw, Inches(0.36), col)
    txt(s, x+Inches(0.1), Inches(1.46), cw-Inches(0.15), Inches(0.28),
        [[(cat, 11, True, WHITE)]], align=PP_ALIGN.CENTER)
    rrect(s, x, Inches(1.82), cw, Inches(3.98), LIGHT, line=col, lw=0.8)
    for j, (name, sub) in enumerate(items):
        txt(s, x+Inches(0.18), Inches(2.02)+j*Inches(1.22), cw-Inches(0.3), Inches(0.38),
            [[(name, 13, True, DARKTX)]])
        txt(s, x+Inches(0.18), Inches(2.42)+j*Inches(1.22), cw-Inches(0.3), Inches(0.38),
            [[(sub, 10.5, False, GREY)]])
rrect(s, Inches(0.55), Inches(5.95), Inches(12.2), Inches(0.7), DARK)
txt(s, Inches(0.78), Inches(6.05), Inches(11.8), Inches(0.5),
    [[("kubectl apply -f catalog.yaml  →  kubectl get aiprovider,aimodel,aisovereigntypolicy,aibudgetpolicy -A", 11, True, GREEN)]],
    align=PP_ALIGN.CENTER)
footer(s)
set_notes(s, """Toute la gouvernance est exprimée en 11 CRDs Kubernetes dans le groupe API aiops.imperium.io.

Le catalogue : AIGateway, AIProvider, AIModel. C'est la déclaration de votre écosystème LLM — qui est là, à quel prix, dans quelle zone.

Coût et budget : AIFinOpsReport pour le reporting, AIBudgetPolicy pour les limites, AIBreakEvenAnalysis pour les recommandations.

Routage : AIRoutingPolicy pour les règles de canary, AIRouteOverride pour les interventions d'urgence, AIChangeRequest pour le circuit d'approbation humaine.

Qualité et souveraineté : AIQualityGate et AISovereigntyPolicy.

Tout est GitOps-ready — vous déclarez l'état désiré, l'opérateur réconcilie en continu.""")

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 15 — 23 MÉTRIQUES
# ═══════════════════════════════════════════════════════════════════════════
s = new_slide()
title_bar(s, "OBSERVABILITÉ", "23 métriques ai_finops_* — Prometheus · Grafana radar 5D", 15)
cols_m = [
    (GREEN,  "Coût & tokens",      ["ai_finops_cost_eur","ai_finops_input_tokens","ai_finops_output_tokens","ai_finops_cost_by_zone_eur","ai_finops_projected_monthly_cost_eur"]),
    (BLUE,   "Budget & économies", ["ai_finops_budget_usage_percent","ai_finops_potential_savings_eur","ai_finops_potential_savings_by_app_eur","ai_finops_cost_saving_eur","ai_finops_recommendations"]),
    (RED,    "Souveraineté",       ["ai_finops_sovereignty_findings","ai_finops_sovereignty_requests","ai_finops_shadow_ai_egress","ai_finops_enforcement_actions","ai_finops_requests"]),
    (ORANGE, "Routing Score (5D)", ["ai_finops_routing_score","ai_finops_cost_score","ai_finops_quality_score","ai_finops_latency_score","ai_finops_reliability_score","ai_finops_sovereignty_score","ai_finops_measured_latency_millis","ai_finops_latency_telemetry_available"]),
]
cw = Inches(2.95)
for i, (col, cat, mets) in enumerate(cols_m):
    x = Inches(0.55)+i*(cw+Inches(0.1))
    rect(s, x, Inches(1.42), cw, Inches(0.36), col)
    txt(s, x+Inches(0.1), Inches(1.46), cw-Inches(0.15), Inches(0.28),
        [[(cat, 11, True, WHITE)]], align=PP_ALIGN.CENTER)
    rrect(s, x, Inches(1.82), cw, Inches(4.52), LIGHT, line=col, lw=0.8)
    txt(s, x+Inches(0.12), Inches(1.94), cw-Inches(0.2), Inches(4.3),
        [[(("\n".join(f"• {m}" for m in mets)), 10.5, False, DARKTX)]], ls=1.2)
rrect(s, Inches(0.55), Inches(6.5), Inches(12.2), Inches(0.65), DARK)
txt(s, Inches(0.78), Inches(6.58), Inches(11.8), Inches(0.45),
    [[("Tous les labels incluent namespace + application → chargeback sans configuration supplémentaire", 12, False, MUTE)]],
    align=PP_ALIGN.CENTER)
footer(s)
set_notes(s, """L'opérateur expose 23 métriques Prometheus sous le préfixe ai_finops_.

Les quatre catégories : coût et tokens — la facturation réelle. Budget et économies — le contrôle des dépenses et les recommandations. Souveraineté — les violations et le shadow AI. Et le routing score — les 5 dimensions du radar Grafana.

Ce qui est important : toutes ces métriques portent les labels namespace et application. Ça veut dire que dans Grafana, vous pouvez filtrer par équipe sans aucune configuration supplémentaire.

Le radar 5D est la visualisation la plus parlante — en un coup d'œil, vous voyez le profil complet de chaque modèle sur les 5 dimensions.""")

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 16 — DÉPLOIEMENT
# ═══════════════════════════════════════════════════════════════════════════
s = new_slide()
title_bar(s, "DÉPLOIEMENT", "3 étapes — Kind local ou AKS production", 16)
steps = [
    (GREEN,  "1.  Opérateur",  ["helm upgrade --install greenops charts/ai-sovereign-finops-operator",
                                  "  --namespace greenops-system --create-namespace",
                                  "# CRDs + Deployment + RBAC + ServiceMonitor opt."]),
    (BLUE,   "2.  Catalogue",  ["kubectl apply -f 06-openai-fr.yaml   # Azure OpenAI France",
                                  "kubectl apply -f 07-openai-us.yaml   # Azure OpenAI US",
                                  "kubectl apply -f 05-mistral-eu.yaml  # Mistral EU Foundry"]),
    (ORANGE, "3.  Politiques", ["kubectl apply -f sovereignty.yaml   # AISovereigntyPolicy",
                                  "kubectl apply -f budgets.yaml        # AIBudgetPolicy × 4",
                                  "# Opérateur réconcilie toutes les 30s"]),
]
for i, (col, title_s, cmds) in enumerate(steps):
    y = Inches(1.42)+i*Inches(1.72)
    rrect(s, Inches(0.55), y, Inches(12.2), Inches(1.62), LIGHT, line=col, lw=1.2)
    txt(s, Inches(0.78), y+Inches(0.1), Inches(11.8), Inches(0.34),
        [[(title_s, 13, True, col)]])
    txt(s, Inches(0.78), y+Inches(0.46), Inches(11.8), Inches(1.05),
        [[(c, 11, False, RGBColor(0x1A,0x2E,0x4A))] for c in cmds], ls=1.12)
rrect(s, Inches(0.55), Inches(6.62), Inches(12.2), Inches(0.52), DARK)
txt(s, Inches(0.78), Inches(6.7), Inches(11.8), Inches(0.35),
    [[("Démo complète 1 commande : ", 12, True, GREEN),
      ("cd automatisation/envoy-aigw  &&  ./deploy.sh up", 12, False, MUTE)]])
footer(s)
set_notes(s, """Le déploiement tient en trois étapes.

Première étape : installer l'opérateur via Helm. Une seule commande, ça crée le namespace, installe les 11 CRDs, le deployment, le RBAC et optionnellement le ServiceMonitor Prometheus.

Deuxième étape : déclarer le catalogue. On applique les fichiers YAML pour chaque provider — OpenAI France, OpenAI US, Mistral EU. Chaque fichier crée les ressources Envoy AI Gateway et les objets opérateur correspondants.

Troisième étape : déclarer les politiques. AISovereigntyPolicy pour la conformité, AIBudgetPolicy pour chaque équipe. L'opérateur commence à réconcilier toutes les 30 secondes.

Pour reproduire la démo complète depuis zéro : une seule commande, deploy.sh up. Elle crée le cluster Kind, construit l'image, installe tout, et démarre les applications de test.""")

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 17 — SYNTHÈSE
# ═══════════════════════════════════════════════════════════════════════════
s = new_slide()
title_bar(s, "SYNTHÈSE", "9 problèmes → 9 solutions → impact mesurable", 17)
rows_s = [
    ("💸", "Facture globale",          "Attribution namespace/app",     "Chargeback exact dès J+1"),
    ("📊", "Budget dépassé sans alerte","AIBudgetPolicy + enforcement",  "Alerte + blocage automatique"),
    ("🔒", "RGPD/AI Act non vérifiés", "AISovereigntyPolicy FR/EU/US",  "Findings critiques temps réel"),
    ("🎯", "Choix modèle intuitif",    "Routing Score 5D + radar",      "Décision sur 5 critères"),
    ("👻", "Shadow AI invisible",       "Tetragon eBPF tcp:443",         "Détection < 30s"),
    ("✅", "Changements sans validation","AIChangeRequest Pending→OK",   "Audit trail complet"),
    ("🔀", "Routage manuel risqué",    "AIRoutingPolicy + Override",    "Rollback en 10 secondes"),
    ("⭐", "Qualité silencieuse",       "AIQualityGate tier minimum",    "Alerte avant l'utilisateur"),
    ("💡", "Savings invisibles",        "AIBreakEvenAnalysis tokens réels","ROI chiffré par équipe"),
]
for i, (ico, prob, sol, imp) in enumerate(rows_s):
    y = Inches(1.4)+i*Inches(0.565)
    bg = LIGHT if i%2==0 else WHITE
    rrect(s, Inches(0.55), y, Inches(12.2), Inches(0.54), bg)
    txt(s, Inches(0.65), y+Inches(0.1), Inches(0.38), Inches(0.35),
        [[(ico, 14, False, DARKTX)]], align=PP_ALIGN.CENTER)
    txt(s, Inches(1.1), y+Inches(0.1), Inches(2.55), Inches(0.35),
        [[(prob, 11, True, RED)]])
    txt(s, Inches(3.78), y+Inches(0.1), Inches(4.4), Inches(0.35),
        [[(sol, 11, False, DARKTX)]])
    txt(s, Inches(8.35), y+Inches(0.1), Inches(4.15), Inches(0.35),
        [[(imp, 11, True, GREEN)]])
rrect(s, Inches(0.55), Inches(6.52), Inches(12.2), Inches(0.62), DARK)
txt(s, Inches(0.78), Inches(6.62), Inches(11.8), Inches(0.42),
    [[("11 CRDs  ·  23 métriques  ·  3 providers réels  ·  Apache 2.0  ·  v0.5.2", 13, True, GREEN)]],
    align=PP_ALIGN.CENTER)
footer(s)
set_notes(s, """En synthèse, 9 problèmes concrets, 9 fonctionnalités, 9 impacts mesurables.

La facture globale → attribution par namespace, chargeback exact dès J+1.
Le budget dépassé → AIBudgetPolicy avec phases et enforcement automatique.
Le RGPD non vérifié → AISovereigntyPolicy avec trois zones en démo live.
Le choix de modèle intuitif → Routing Score 5D visible dans le radar Grafana.
Le Shadow AI invisible → Tetragon eBPF, détection en moins de 30 secondes.
Les changements sans validation → AIChangeRequest avec circuit d'approbation.
Le routage manuel → AIRoutingPolicy et AIRouteOverride, rollback en 10 secondes.
La qualité silencieuse → AIQualityGate avec tier minimum configurable.
Les économies invisibles → AIBreakEvenAnalysis sur tokens réels.

Tout ça en 11 CRDs, 23 métriques, Apache 2.0.""")

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 18 — LIMITES
# ═══════════════════════════════════════════════════════════════════════════
s = new_slide()
title_bar(s, "HONNÊTETÉ", "Limites actuelles — assumées", 18)
limits = [
    (ORANGE, "Latence ≠ TTFT",            "Score basé sur durée gateway (P50). Time-to-first-token non séparé."),
    (ORANGE, "Quality score = tier statique","Déclaré dans l'AIModel, pas évaluation LLM-as-a-judge temps réel."),
    (ORANGE, "Enforcement sur Envoy seul", "Kong, NGINX, Istio non supportés dans cette version."),
    (ORANGE, "Tetragon = privilèges noyau","Requiert CAP_BPF. Désactivable sans impact sur le plan gateway."),
]
for i, (col, h, b) in enumerate(limits):
    y = Inches(1.42)+i*Inches(1.35)
    rrect(s, Inches(0.55), y, Inches(12.2), Inches(1.25), ORGBG, line=col, lw=1.0)
    txt(s, Inches(0.78), y+Inches(0.1), Inches(11.8), Inches(0.36),
        [[(h, 13, True, col)]])
    txt(s, Inches(0.78), y+Inches(0.5), Inches(11.8), Inches(0.65),
        [[(b, 12.5, False, DARKTX)]], ls=1.1)
footer(s)
set_notes(s, """On termine avec honnêteté sur les limites actuelles.

Le score de latence est calculé sur la durée totale mesurée par le gateway — le time-to-first-token, qui est souvent la métrique la plus importante pour l'expérience utilisateur, n'est pas encore isolé.

Le quality score est statique — il reflète le tier déclaré dans l'AIModel, pas une évaluation automatique de la qualité des réponses. C'est une limite connue qu'on assume.

L'enforcement ne fonctionne qu'avec Envoy AI Gateway dans cette version. Kong, NGINX et Istio ne sont pas encore supportés.

Tetragon nécessite des privilèges noyau. Dans des environnements très restrictifs, on peut désactiver le plan eBPF — le plan gateway reste totalement fonctionnel.

Ces limites sont documentées dans le README. On préfère être transparents plutôt que de surprendre en production.""")

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 19 — FIN
# ═══════════════════════════════════════════════════════════════════════════
s = new_slide()
rect(s, 0, 0, SW, SH, DARK)
rect(s, 0, 0, Inches(0.28), SH, GREEN)
rect(s, 0, SH-Pt(4), SW, Pt(4), GREEN)
txt(s, Inches(0.9), Inches(1.5), Inches(11.5), Inches(0.45),
    [[("MERCI", 14, True, GREEN)]])
txt(s, Inches(0.9), Inches(2.1), Inches(11.5), Inches(0.9),
    [[("Questions ?", 52, True, WHITE)]])
txt(s, Inches(0.9), Inches(3.4), Inches(11.6), Inches(0.5),
    [[("Démo live → ", 15, True, GREEN),
      ("http://localhost:3000", 15, False, MUTE),
      ("  ·  Grafana radar + Shadow AI + Budget", 13, False, GREY)]])
for i, (k, v) in enumerate([
    ("GitHub",    "github.com/ihsenalaya/ai-sovereign-finops-operator"),
    ("Licence",   "Apache 2.0  ·  v0.5.2"),
    ("11 CRDs",   "23 métriques  ·  3 providers réels  ·  ./deploy.sh up"),
    ("Contact",   "articlouet@gmail.com"),
]):
    txt(s, Inches(0.9), Inches(4.45)+i*Inches(0.52), Inches(11.6), Inches(0.44),
        [[(k+"  :  ", 13, True, GREEN), (v, 13, False, MUTE)]])
set_notes(s, """Merci pour votre attention.

Le projet est open source sous licence Apache 2.0, disponible sur GitHub. La démo tourne en ce moment sur le cluster Kind local — Grafana est accessible sur localhost:3000.

Pour reproduire la démo chez vous : cloner le repo, lancer deploy.sh up, et vous avez l'environnement complet en moins de 10 minutes.

Je suis disponible pour les questions.""")

# ── OUTPUT ───────────────────────────────────────────────────────────────
out = "AI-Sovereign-FinOps-Presentation.pptx"
prs.save(out)
print(f"✅  {out}  —  {len(prs.slides)} slides  (notes présentateur sur chaque slide)")
